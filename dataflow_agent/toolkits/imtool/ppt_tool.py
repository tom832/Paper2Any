# -*- coding: utf-8 -*-
"""
ppt_tool

本模块将一组顺序图片通过 PaddleOCR 识别为可编辑文本，自动分析行高、版式和颜色，生成带“干净底图+覆盖文字框”的 PPTX，并可选同时导出 PDF；
内部提供多种参数控制背景 inpaint 强度（INPAINT_METHOD/INPAINT_RADIUS、SIMPLE_BG_VAR_THRESH、MASK_DILATE_ITER、USE_ADAPTIVE_MASK）、
OCR 分辨率与锐化（UPSCALE_LONG_SIDE_TO、UPSCALE_INTERP、ENABLE_SHARPEN、SHARPEN_AMOUNT）、
文本过滤阈值（DROP_SCORE）以及字号放大与标题/副标题对正文的比例（BASE_BODY_PT、FONT_SCALE_FACTOR、TITLE_RATIO_*/SUBTITLE_RATIO_*/BODY_RATIO_*），
并通过 ADD_BACKGROUND_IMAGE / CLEAN_BACKGROUND / EXTRACT_TEXT_COLOR 控制是否叠加背景图片、是否抠掉原文字、是否按原图估计文字颜色，从而在“还原视觉效果”与“可编辑性/美观度”和运行性能之间做平衡。

功能概述：
- 从指定目录按自然顺序读取图片
- 生成包含所有图片页的 PDF 文件
- 使用 PaddleOCR 对页面进行 OCR，识别文本行
- 基于识别结果自动估计字体大小和颜色，将文本叠加到 PPTX 中
- 可选地对原始页面进行 inpaint，生成“去文字的干净底图”作为 PPT 背景

典型用法：
- 在 DataFlow-Agent 的图像处理流程中，作为从图片页到可编辑 PPT 文稿的后处理工具
- 也可在其它组件或脚本中通过对外函数直接调用
"""

# 函数一览：
# natural_key(s): 生成用于文件名“自然排序”的 key，将数字部分按整数比较。
# list_images_in_dir(d): 按自然顺序列出目录中的所有图片文件路径。
# read_bgr(path): 以兼容非 ASCII 路径的方式读取图片，并返回标准 BGR uint8 格式。
# debug_dump(img, tag): 将中间图像写入调试目录并记录基础统计信息。
# images_to_pdf(image_paths, output_pdf_path): 将一组图片顺序导出为单个 PDF 文件。
# pdf_to_images(pdf_path, out_dir, dpi): 将 PDF 每一页按指定分辨率渲染为 PNG，并返回图片路径列表。
# upscale_if_needed(bgr, long_side_to, interp): 若分辨率偏低则按长边放大图像，返回放大后图像和缩放比例。
# sharpen(bgr, amount): 使用“反锐化掩模”方式对图像进行轻度锐化。
# preprocess_for_ocr(bgr): 对整页图像做放大与可选锐化，生成适合 OCR 的版本及缩放比例。
# is_cjk(s): 判断字符串中是否包含 CJK（中日韩）字符。
# iou(a, b): 计算两个矩形框的交并比（IoU）。
# merge_lines(lines, y_tol, x_gap): 将 OCR 的短行/单词按行方向与间距合并成句级文本行。
# text_score(lines): 根据字符数量、平均置信度及是否含 CJK，估计一组文本行的整体得分。
# paddle_ocr(bgr, drop_score): 调用 PaddleOCR 对整页 BGR 图像做 OCR，并按置信度阈值过滤结果。
# paddle_ocr_page_with_layout(img_path): 对单页图片做预处理 + OCR + 行合并 + 行高/背景色估计并返回布局信息。
# extract_text_color(bgr, bbox, bg_color): 从给定文字区域估计主文字颜色，尽量排除接近背景的颜色。
# estimate_background_color(bgr, lines): 用文字 mask 反选背景区域，估计页面主背景颜色。
# px_to_emu(px, emu_per_px): 将像素值按给定比例转换为 PPT 使用的 EMU 单位。
# analyze_line_heights(lines): 统计 OCR 行框高度分布，估计正文行高的中位数。
# classify_line_role(bbox, img_h_px, body_h_px): 根据行高与垂直位置粗略区分标题、副标题和正文。
# estimate_font_pt(bbox, img_h_px, body_h_px, slide_h_in): 依据原图行高比例估计在 PPT 中的字号大小。
# add_background(slide, bgr, slide_w_emu, slide_h_emu, tmp_path): 将整页背景图添加到 PPT 幻灯片并删除临时文件。
# build_text_mask_from_lines(bgr, lines): 基于 OCR 行框生成粗略的文字区域二值 mask。
# build_adaptive_mask(bgr, lines): 结合局部对比度与自适应阈值，生成更精细的文字主 mask。
# is_simple_background_region(bgr, mask): 判断文字区域邻域背景是否近似纯色（方差较小）。
# fill_with_neighbor(bgr, mask): 对复杂背景的文字区域先用邻域像素进行粗填充，缓解 inpaint 伪影。
# make_clean_background(bgr, lines): 基于文字 mask 和 inpaint 生成“去文字的干净底图”。
# ocr_images_to_ppt(image_paths, output_pptx, add_background_image, clean_background, use_text_color): 将图片序列通过 OCR 转成带背景与覆盖文本框的可编辑 PPT。
# images_to_pdf_and_ppt(image_paths, output_pdf_path, output_pptx_path, add_background_image, clean_background, extract_text_color): 将给定图片列表一站式转换为 PDF 和 PPTX 并返回路径。
# convert_images_dir_to_pdf_and_ppt(input_dir, output_pdf_path, output_pptx_path, add_background_image, clean_background, extract_text_color): 从图片目录读取图片并生成对应的 PDF + PPTX。
# convert_images_dir_to_pdf_and_ppt_api(input_dir, output_pdf_path, output_pptx_path, api_url, api_key, model, use_api_inpaint, add_background_image, clean_background, use_text_color): 异步版本的目录转 PDF/PPTX，优先使用图像编辑 API 做 inpainting，失败时回退到本地 inpaint。

import os
import re
from typing import Sequence, Optional, Dict, Any, List, Tuple
import requests
import random
from collections import Counter

import fitz  # PyMuPDF
from pathlib import Path

import numpy as np
from PIL import Image
import cv2
from paddleocr import PaddleOCR

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from dataflow_agent.utils import get_project_root
from dataflow_agent.logger import get_logger
from typing import Union

log = get_logger(__name__)

# ----------------------------
# Config (默认配置，可通过对外函数参数进行部分覆盖)
# ----------------------------
ADD_BACKGROUND_IMAGE = True
CLEAN_BACKGROUND = True  # 是否尝试抠掉文字、生成无字底图再叠加OCR文字
EXTRACT_TEXT_COLOR = True  # 是否提取原图文字颜色
INPAINT_METHOD = cv2.INPAINT_TELEA  # or cv2.INPAINT_NS
INPAINT_RADIUS = 7  # 增大修复半径（从3提高到7）
SIMPLE_BG_VAR_THRESH = 50.0  # 放宽阈值（从12提高到50）
MASK_DILATE_ITER = 2  # 增加膨胀次数（从1提高到2）
USE_ADAPTIVE_MASK = True  # 使用自适应mask生成

# 输出PPT比例（16:9）
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Debug：落盘你送入 OCR 的图片，方便肉眼确认内容/分辨率/是否被处理坏
DEBUG_DUMP_FIRST_N = 2
DEBUG_DIR = f"{get_project_root()}/tests/debug_frames"

# ---------- 核心修复：低分辨率页面的 OCR 前增强 ----------
UPSCALE_LONG_SIDE_TO = 2200  # 建议 2000~3200，越大越慢
UPSCALE_INTERP = cv2.INTER_CUBIC
ENABLE_SHARPEN = True  # 轻度锐化，提升边缘对比度
SHARPEN_AMOUNT = 0.8  # 0.0~1.5 之间

# 识别过滤阈值
DROP_SCORE = 30  # PaddleOCR 的score是 0-1，这里统一乘100后按0-100过滤

# 字号优化配置（放大整体字号，明显拉开标题对正文的差距）
BASE_BODY_PT = 16.0  # 正文基准字号
FONT_SCALE_FACTOR = 1.0  # 全局字号缩放因子
TITLE_RATIO_MIN = 2.0  # 标题最小倍率
TITLE_RATIO_MAX = 3.5  # 标题最大倍率
SUBTITLE_RATIO_MIN = 1.4  # 副标题最小倍率
SUBTITLE_RATIO_MAX = 2.0  # 副标题最大倍率
BODY_RATIO_MIN = 0.9  # 正文最小倍率
BODY_RATIO_MAX = 1.1  # 正文最大倍率

# PaddleOCR 配置（全局只初始化一次）
PADDLE_OCR = PaddleOCR(
    use_angle_cls=True,  # 角度分类，处理横竖混排
    lang="ch",  # 中文 + 英文
)

# ----------------------------
# Font Size Clustering
# ----------------------------

class FontSizeClustering:
    """
    自适应字号聚类器：将连续的、有噪声的字号估算值，映射到 K 个离散的"标准字号"。
    支持全局聚类（全文档统一）或单页聚类。
    依赖 sklearn.cluster.KMeans，如果缺失则回退到简单的众数/分位数策略。
    """
    def __init__(self, n_clusters: int = 4, merge_tol: float = 2.0):
        self.n_clusters = n_clusters
        self.merge_tol = merge_tol
        self.centroids = []
        self.has_sklearn = False
        try:
            from sklearn.cluster import KMeans
            self._KMeans = KMeans
            self.has_sklearn = True
        except ImportError:
            pass

    def fit(self, font_sizes: List[float]) -> "FontSizeClustering":
        """
        输入原始字号列表（pt），计算聚类中心。
        """
        # 过滤无效值
        data = [x for x in font_sizes if x > 0]
        if not data:
            self.centroids = [12.0] # 默认回退
            return self

        # 1. 如果数据量太少，直接用原始值（去重排序）
        if len(data) < self.n_clusters:
            self.centroids = sorted(list(set(data)))
            return self

        # 2. 如果没有 sklearn，回退到简单的直方图统计（取前K个高频值）
        if not self.has_sklearn:
            # 简单统计：取出现频率最高的 K 个，或者简单的分位数
            # 这里用频率统计更符合"标准字号"的直觉
            counts = Counter([round(x) for x in data])
            top_k = counts.most_common(self.n_clusters)
            self.centroids = sorted([float(x[0]) for x in top_k])
            return self

        # 3. K-Means 聚类
        import numpy as np
        X = np.array(data).reshape(-1, 1)
        
        # 动态调整 K：不能超过样本唯一值的数量
        n_unique = len(set([round(x, 1) for x in data]))
        real_k = min(self.n_clusters, n_unique)
        
        kmeans = self._KMeans(n_clusters=real_k, n_init=10, random_state=42)
        kmeans.fit(X)
        centers = sorted(kmeans.cluster_centers_.flatten())

        # 4. 后处理：合并过近的中心 (Merge close centers)
        merged_centers = []
        if centers:
            curr = centers[0]
            for next_c in centers[1:]:
                if (next_c - curr) < self.merge_tol:
                    # 距离太近，合并（取平均）
                    curr = (curr + next_c) / 2.0
                else:
                    merged_centers.append(curr)
                    curr = next_c
            merged_centers.append(curr)
        
        # 圆整到 0.5 pt
        self.centroids = [round(c * 2) / 2.0 for c in merged_centers]
        log.info(f"[FontSizeClustering] Fitted centroids: {self.centroids}")
        return self

    def map(self, pt: float) -> float:
        """
        将原始字号映射到最近的中心。
        """
        if not self.centroids:
            return pt
        
        # 找最近邻
        closest = min(self.centroids, key=lambda c: abs(c - pt))
        return closest


# ----------------------------
# IO helpers
# ----------------------------


def natural_key(s: str):
    return [int(t) if t.isdigit() else t.lower() for t in re.split(r"(\d+)", s)]


def list_images_in_dir(d: str) -> List[str]:
    """
    按自然顺序列出目录中所有图片文件路径。
    """
    exts = (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff")
    files = [f for f in os.listdir(d) if f.lower().endswith(exts)]
    files.sort(key=natural_key)
    return [os.path.join(d, f) for f in files]


def read_bgr(path: str) -> np.ndarray:
    """
    Robust image reader:
    - supports non-ascii paths (np.fromfile + imdecode)
    - returns BGR uint8 HxWx3
    """
    data = np.fromfile(path, dtype=np.uint8)
    img = cv2.imdecode(data, cv2.IMREAD_UNCHANGED)
    if img is None:
        img = cv2.imread(path, cv2.IMREAD_UNCHANGED)
    if img is None:
        raise ValueError(f"Failed to read image: {path}")

    # Normalize to BGR uint8
    if img.ndim == 2:
        img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)
    elif img.ndim == 3 and img.shape[2] == 4:
        img = cv2.cvtColor(img, cv2.COLOR_BGRA2BGR)

    if img.dtype != np.uint8:
        img = np.clip(img, 0, 255).astype(np.uint8)

    return img


def debug_dump(img: np.ndarray, tag: str = "dbg") -> None:
    """
    将中间图像写入 DEBUG_DIR 方便调试。
    """
    os.makedirs(DEBUG_DIR, exist_ok=True)

    log.info(f"{tag} type: {type(img)}")
    if isinstance(img, np.ndarray):
        log.info(
            f"{tag} shape: {img.shape}, dtype: {img.dtype}, "
            f"min/max: {int(img.min())}/{int(img.max())}"
        )

    out_path = os.path.join(DEBUG_DIR, f"{tag}.png")
    ok = cv2.imwrite(out_path, img)
    log.info(f"{tag} saved: {out_path}, ok: {ok}")


# ----------------------------
# PDF / Page helpers
# ----------------------------


def images_to_pdf(image_paths: Sequence[str], output_pdf_path: str) -> str:
    """
    将一组图片导出为单个 PDF 文件。
    """
    imgs: List[Image.Image] = []
    for p in image_paths:
        im = Image.open(p)
        if im.mode != "RGB":
            im = im.convert("RGB")
        imgs.append(im)
    if not imgs:
        raise ValueError("No images for PDF.")
    imgs[0].save(output_pdf_path, save_all=True, append_images=imgs[1:])
    return output_pdf_path


def pdf_to_images(pdf_path: str, out_dir: str, dpi: int = 220) -> List[str]:
    """
    将 PDF 每一页渲染为 PNG 图片，返回图片路径列表（按页码顺序）。

    参数
    ----
    pdf_path:
        输入 PDF 文件路径。
    out_dir:
        输出图片所在目录，不存在会自动创建。
    dpi:
        渲染分辨率（每英寸像素数），默认 220。

    返回
    ----
    List[str]
        按页码顺序排列的 PNG 图片绝对路径列表。
    """
    doc = fitz.open(pdf_path)
    out_dir_path = Path(out_dir)
    out_dir_path.mkdir(parents=True, exist_ok=True)

    image_paths: List[str] = []
    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)

        img_path = out_dir_path / f"page_{page_index + 1:03d}.png"
        pix.save(str(img_path))
        image_paths.append(str(img_path))

    doc.close()
    log.info(f"[pdf_to_images] rendered {len(image_paths)} pages from {pdf_path}")
    return image_paths


# ----------------------------
# Preprocess
# ----------------------------


def upscale_if_needed(
    bgr: np.ndarray,
    long_side_to: int = UPSCALE_LONG_SIDE_TO,
    interp: int = UPSCALE_INTERP,
):
    h, w = bgr.shape[:2]
    long_side = max(h, w)
    if long_side >= long_side_to:
        return bgr, 1.0

    scale = long_side_to / float(long_side)
    new_w = int(round(w * scale))
    new_h = int(round(h * scale))
    up = cv2.resize(bgr, (new_w, new_h), interpolation=interp)
    return up, scale


def sharpen(bgr: np.ndarray, amount: float = SHARPEN_AMOUNT) -> np.ndarray:
    """
    Unsharp mask style: sharpen = img*(1+a) - blur*a
    """
    if amount <= 0:
        return bgr
    blur = cv2.GaussianBlur(bgr, (0, 0), sigmaX=1.2, sigmaY=1.2)
    out = cv2.addWeighted(bgr, 1.0 + amount, blur, -amount, 0)
    return np.clip(out, 0, 255).astype(np.uint8)


def preprocess_for_ocr(bgr: np.ndarray):
    """
    Make a "det-friendly" version of the page:
    - upscale to a reasonable working resolution
    - optional sharpen
    """
    up, scale = upscale_if_needed(bgr)
    if ENABLE_SHARPEN:
        up = sharpen(up, amount=SHARPEN_AMOUNT)
    return up, scale


# ----------------------------
# OCR helpers
# ----------------------------


def is_cjk(s: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in s)


def iou(a, b) -> float:
    ax1, ay1, ax2, ay2 = a
    bx1, by1, bx2, by2 = b
    ix1, iy1 = max(ax1, bx1), max(ay1, by1)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
    inter = iw * ih
    if inter <= 0:
        return 0.0
    area_a = (ax2 - ax1) * (ay2 - ay1)
    area_b = (bx2 - bx1) * (by2 - by1)
    return inter / (area_a + area_b - inter + 1e-6)


def merge_lines(
    lines: Sequence[Tuple[Sequence[float], str, float]], y_tol: int = 12, x_gap: int = 18
):
    """
    将OCR的word/短行合并成句子级别的行
    """
    if not lines:
        return []
    lines = sorted(lines, key=lambda x: (x[0][1], x[0][0]))

    def union(b1, b2):
        return [
            min(b1[0], b2[0]),
            min(b1[1], b2[1]),
            max(b1[2], b2[2]),
            max(b1[3], b2[3]),
        ]

    merged = []
    cur_bbox, cur_text, cur_conf_sum, cur_n = (
        lines[0][0],
        lines[0][1],
        lines[0][2],
        1,
    )

    for bbox, text, conf in lines[1:]:
        cy1 = (cur_bbox[1] + cur_bbox[3]) / 2
        cy2 = (bbox[1] + bbox[3]) / 2
        same_line = abs(cy1 - cy2) <= y_tol
        near_x = (bbox[0] - cur_bbox[2]) <= x_gap

        if same_line and near_x:
            cur_bbox = union(cur_bbox, bbox)
            if (not is_cjk(cur_text)) and (not is_cjk(text)):
                cur_text = (cur_text + " " + text).strip()
            else:
                cur_text = (cur_text + text).strip()
            cur_conf_sum += conf
            cur_n += 1
        else:
            merged.append((cur_bbox, cur_text, cur_conf_sum / cur_n))
            cur_bbox, cur_text, cur_conf_sum, cur_n = bbox, text, conf, 1

    merged.append((cur_bbox, cur_text, cur_conf_sum / cur_n))
    return merged


def text_score(lines) -> float:
    if not lines:
        return 0.0
    total_chars = sum(len(t) for (_, t, _) in lines)
    avg_conf = sum(conf for (_, _, conf) in lines) / max(1, len(lines))
    cjk_bonus = 1.1 if any(is_cjk(t) for (_, t, _) in lines) else 1.0
    return total_chars * (avg_conf / 100.0) * cjk_bonus  # normalize confidence to 0-1


def paddle_ocr(bgr: np.ndarray, drop_score: int = DROP_SCORE):
    """
    使用 PaddleOCR 识别整页图片
    返回格式：[(bbox, text, confidence), ...]
    bbox: [x1, y1, x2, y2]
    注意：这里直接在 BGR 图上跑，PaddleOCR 内部会处理颜色空间。
    """
    h, w = bgr.shape[:2]

    # ocr_result: List[List[ [box, (text, score)], ... ]]
    ocr_result = PADDLE_OCR.ocr(bgr, cls=True)
    lines = []

    if not ocr_result:
        return lines

    # 通常一页对应 ocr_result[0]
    for line in ocr_result[0]:
        box, (text, score) = line
        if not text:
            continue
        if score * 100.0 < drop_score:
            continue

        # box 是四点多边形：[ [x1,y1], [x2,y2], [x3,y3], [x4,y4] ]
        xs = [p[0] for p in box]
        ys = [p[1] for p in box]
        x1, y1, x2, y2 = min(xs), min(ys), max(xs), max(ys)

        # 边界裁剪
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            continue

        bbox = [float(x1), float(y1), float(x2), float(y2)]
        # 保持 0-100 置信度范围，和原 Tesseract 逻辑兼容
        lines.append((bbox, text.strip(), float(score * 100.0)))

    return lines


def paddle_ocr_page_with_layout(img_path: str) -> Dict[str, Any]:
    """
    对单页图片执行：
    - 读取 + 预处理(放大 + 锐化)
    - PaddleOCR 识别
    - 坐标从 OCR 分辨率映射回原图
    - 行合并
    - 正文行高估计
    - 背景颜色估计

    返回:
    {
        "image_size": (w, h),
        "lines": [(bbox, text, conf), ...],  # bbox 为原图像素坐标
        "body_h_px": float 或 None,
        "bg_color": (r,g,b) 或 None,
    }
    """
    bgr = read_bgr(img_path)
    h0, w0 = bgr.shape[:2]

    # 预处理
    ocr_img, scale = preprocess_for_ocr(bgr)
    h1, w1 = ocr_img.shape[:2]

    log.info(f"[paddle_ocr_page_with_layout] {os.path.basename(img_path)} up-scale={scale:.3f}")

    # OCR
    raw_lines = paddle_ocr(ocr_img)

    # 映射回原图像素坐标
    if raw_lines and (w1 != w0 or h1 != h0):
        sx = w0 / float(w1)
        sy = h0 / float(h1)
        raw_lines = [
            ([b[0] * sx, b[1] * sy, b[2] * sx, b[3] * sy], t, c)
            for (b, t, c) in raw_lines
        ]

    # 合并行
    y_tol = max(12, int(h0 * 0.008))
    x_gap = max(18, int(w0 * 0.01))
    lines = merge_lines(raw_lines, y_tol=y_tol, x_gap=x_gap)

    # 正文行高估计
    body_h_px = analyze_line_heights(lines)

    if not lines:
        log.warning(f"[paddle_ocr_page_with_layout] no text detected: {img_path}")
        bg_color = None
    else:
        log.info(
            f"[paddle_ocr_page_with_layout] detected {len(lines)} text boxes, body_h_px={body_h_px}"
        )
        bg_color = estimate_background_color(bgr, lines) if EXTRACT_TEXT_COLOR else None

    return {
        "image_size": (w0, h0),
        "lines": lines,
        "body_h_px": body_h_px,
        "bg_color": bg_color,
    }


def paddle_ocr_page_with_layout_server(
    img_path: str,
    server_urls: Union[str, List[str]],
) -> Dict[str, Any]:
    """
    对单页图片执行远程 OCR 处理。

    参数:
        img_path: 图片路径
        server_urls: OCR 服务器 URL 或 URL 列表

    返回:
        同 paddle_ocr_page_with_layout
    """
    if isinstance(server_urls, str):
        urls = [server_urls]
    else:
        urls = list(server_urls)
    
    if not urls:
        raise ValueError("No server URLs provided")
    
    base_url = random.choice(urls)
    api_url = f"{base_url.rstrip('/')}/predict"
    
    abs_img_path = os.path.abspath(img_path)
    payload = {"image_path": abs_img_path}
    
    try:
        response = requests.post(api_url, json=payload, timeout=300)
        response.raise_for_status()
        data = response.json()
        
        # Transform lines back to tuples if needed, though list is fine
        # lines: [[bbox, text, conf], ...] -> [(bbox, text, conf), ...]
        lines = []
        for line_obj in data.get("lines", []):
            lines.append((
                line_obj.get("bbox"),
                line_obj.get("text"),
                line_obj.get("conf")
            ))
            
        return {
            "image_size": tuple(data.get("image_size", [0, 0])),
            "lines": lines,
            "body_h_px": data.get("body_h_px"),
            "bg_color": tuple(data.get("bg_color")) if data.get("bg_color") else None
        }
        
    except Exception as e:
        raise RuntimeError(f"Failed to call OCR server at {api_url}: {e}")


# ----------------------------
# Color extraction
# ----------------------------


def extract_text_color(
    bgr: np.ndarray, bbox, bg_color=None
) -> Tuple[int, int, int]:
    """
    从文字区域提取主色调
    返回 (r, g, b) 元组
    """
    x1, y1, x2, y2 = [int(round(v)) for v in bbox]
    h, w = bgr.shape[:2]

    # 边界检查
    x1 = max(0, min(w - 1, x1))
    x2 = max(0, min(w, x2))
    y1 = max(0, min(h - 1, y1))
    y2 = max(0, min(h, y2))

    if x2 <= x1 or y2 <= y1:
        return (0, 0, 0)  # 默认黑色

    # 提取文字区域
    region = bgr[y1:y2, x1:x2]
    if region.size == 0:
        return (0, 0, 0)

    # 转换为RGB
    region_rgb = cv2.cvtColor(region, cv2.COLOR_BGR2RGB)
    pixels = region_rgb.reshape(-1, 3)

    # 如果像素太少，直接返回中位数颜色
    if len(pixels) < 10:
        median_color = np.median(pixels, axis=0).astype(int)
        return tuple(int(x) for x in median_color)

    # 使用K-means聚类找出主色调（2-3个聚类）
    try:
        from sklearn.cluster import KMeans

        n_clusters = min(3, len(pixels))
        kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
        kmeans.fit(pixels)

        # 获取聚类中心和每个聚类的像素数
        centers = kmeans.cluster_centers_
        labels = kmeans.labels_
        counts = np.bincount(labels)

        # 如果提供了背景色，排除接近背景色的聚类
        if bg_color is not None:
            bg_array = np.array(bg_color)
            valid_centers = []
            valid_counts = []

            for i, center in enumerate(centers):
                # 计算与背景色的距离
                dist = np.linalg.norm(center - bg_array)
                if dist > 30:  # 距离阈值
                    valid_centers.append(center)
                    valid_counts.append(counts[i])

            if valid_centers:
                centers = np.array(valid_centers)
                counts = np.array(valid_counts)

        # 选择出现频率最高的颜色
        dominant_idx = np.argmax(counts)
        dominant_color = centers[dominant_idx].astype(int)

        return tuple(int(x) for x in dominant_color)
    except Exception:
        # 如果sklearn不可用或出错，使用简单的中位数方法
        median_color = np.median(pixels, axis=0).astype(int)
        return tuple(int(x) for x in median_color)


def estimate_background_color(bgr: np.ndarray, lines):
    """
    估计背景主色调，用于颜色提取时排除背景
    """
    h, w = bgr.shape[:2]

    # 创建文字mask
    mask = np.ones((h, w), dtype=np.uint8) * 255
    for bbox, _, _ in lines:
        x1, y1, x2, y2 = [int(round(v)) for v in bbox]
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 > x1 and y2 > y1:
            mask[y1:y2, x1:x2] = 0

    # 提取背景区域像素
    bg_pixels = bgr[mask > 0]
    if bg_pixels.size == 0:
        return None

    # 转换为RGB并计算中位数
    bg_rgb = cv2.cvtColor(
        bg_pixels.reshape(-1, 1, 3), cv2.COLOR_BGR2RGB
    ).reshape(-1, 3)
    median_bg = np.median(bg_rgb, axis=0).astype(int)

    return tuple(int(x) for x in median_bg)


# ----------------------------
# PPT helpers
# ----------------------------


def px_to_emu(px: float, emu_per_px: float) -> int:
    return int(px * emu_per_px)


def analyze_line_heights(lines) -> Optional[float]:
    """
    统计行高分布，估计"正文行高"
    """
    if not lines:
        return None
    hs = [max(1, b[3] - b[1]) for (b, _, _) in lines]
    return float(np.median(hs))


def classify_line_role(bbox, img_h_px: int, body_h_px: Optional[float]) -> str:
    """
    大致区分：title / subtitle / body
    """
    x1, y1, x2, y2 = bbox
    h = max(1, y2 - y1)
    if body_h_px is None or body_h_px <= 0:
        return "body"
    ratio = h / float(body_h_px)

    # 位置辅助：靠近页面顶部 + 较高行
    y_center = (y1 + y2) / 2.0
    top_region = img_h_px * 0.3

    if ratio > 1.7 and y_center < top_region:
        return "title"
    if ratio > 1.3:
        return "subtitle"
    return "body"


def estimate_font_pt(
    bbox, img_h_px: int, body_h_px: Optional[float], slide_h_in: float = SLIDE_H_IN
):
    """
    根据行高按比例估计字号（改进版：移除硬编码倍率限制）
    
    核心思路：
    1. 计算原图中该行的像素高度
    2. 按比例映射到PPT的点数(pt)
    3. 不再强制限制标题/副标题的倍率范围
    """
    x1, y1, x2, y2 = bbox
    h_px = max(1, y2 - y1)
    
    # 方法：将像素高度按图片高度比例转换为PPT点数
    # PPT高度 = 7.5英寸 = 540pt (1英寸=72pt)
    slide_h_pt = slide_h_in * 72.0
    
    # 行高占图片高度的比例
    height_ratio = h_px / float(img_h_px)
    
    # 映射到PPT点数，乘以0.7是经验系数（因为行高通常大于字号）
    pt = slide_h_pt * height_ratio * 0.7
    
    # 应用全局缩放因子
    pt *= FONT_SCALE_FACTOR
    
    # 只做合理范围限制，不再强制角色倍率
    return Pt(max(8, min(96, pt)))


def add_background(
    slide, bgr: np.ndarray, slide_w_emu: int, slide_h_emu: int, tmp_path: str
) -> None:
    rgb = cv2.cvtColor(bgr, cv2.COLOR_BGR2RGB)
    pil = Image.fromarray(rgb)
    pil.save(tmp_path)
    slide.shapes.add_picture(tmp_path, 0, 0, width=slide_w_emu, height=slide_h_emu)
    os.remove(tmp_path)


def build_text_mask_from_lines(bgr: np.ndarray, lines) -> np.ndarray:
    """
    根据OCR行框生成初始mask（粗略矩形）
    """
    h, w = bgr.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)

    for bbox, text, conf in lines:
        x1, y1, x2, y2 = [int(round(v)) for v in bbox]
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            continue
        mask[y1:y2, x1:x2] = 255

    if MASK_DILATE_ITER > 0:
        kernel = np.ones((3, 3), np.uint8)
        mask = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)

    return mask


def build_adaptive_mask(bgr: np.ndarray, lines) -> np.ndarray:
    """
    使用自适应方法生成更精细的文字主mask
    结合OCR bbox和实际文字形状（内部边缘 + 阈值）
    """
    h, w = bgr.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)

    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)

    for bbox, text, conf in lines:
        x1, y1, x2, y2 = [int(round(v)) for v in bbox]
        x1 = max(0, min(w - 1, x1))
        x2 = max(0, min(w, x2))
        y1 = max(0, min(h - 1, y1))
        y2 = max(0, min(h, y2))
        if x2 <= x1 or y2 <= y1:
            continue

        region = gray[y1:y2, x1:x2]
        if region.size == 0:
            continue

        try:
            # 先看局部对比度
            if np.var(region) < 100:
                # 对比度很低，优先用 Canny 边缘找笔画
                edges = cv2.Canny(region, 50, 150)
                kernel = np.ones((2, 2), np.uint8)
                binary = cv2.dilate(edges, kernel, iterations=1)
            else:
                # 对比度正常，用阈值法
                if region.shape[0] < 20 or region.shape[1] < 20:
                    _, binary = cv2.threshold(
                        region, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU
                    )
                    binary = 255 - binary  # 反色：文字为白
                else:
                    binary = cv2.adaptiveThreshold(
                        region,
                        255,
                        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                        cv2.THRESH_BINARY_INV,
                        11,
                        2,
                    )
                kernel = np.ones((2, 2), np.uint8)
                binary = cv2.morphologyEx(binary, cv2.MORPH_CLOSE, kernel)

            mask[y1:y2, x1:x2] = cv2.bitwise_or(mask[y1:y2, x1:x2], binary)
        except Exception:
            mask[y1:y2, x1:x2] = 255

    if MASK_DILATE_ITER > 0:
        kernel = np.ones((3, 3), np.uint8)
        mask = cv2.dilate(mask, kernel, iterations=MASK_DILATE_ITER)

    return mask


def is_simple_background_region(bgr: np.ndarray, mask: np.ndarray) -> bool:
    """
    简单判定：mask 区域附近背景是否接近纯色（方差较小）
    """
    gray = cv2.cvtColor(bgr, cv2.COLOR_BGR2GRAY)
    # 扩大一点范围取邻域
    dilated = cv2.dilate((mask > 0).astype(np.uint8), np.ones((5, 5), np.uint8), iterations=1)
    region = gray[dilated > 0]
    if region.size == 0:
        return False
    var = float(np.var(region))
    return var < SIMPLE_BG_VAR_THRESH


def fill_with_neighbor(bgr: np.ndarray, mask: np.ndarray) -> np.ndarray:
    """
    对复杂背景时，优先用邻域像素粗略填充，再交给 inpaint 做平滑，
    避免 NS/TELEA 在大块区域产生奇怪纹理。
    """
    result = bgr.copy()
    h, w = mask.shape
    for y in range(h):
        xs = np.where(mask[y] > 0)[0]
        if len(xs) == 0:
            continue
        x_min, x_max = xs[0], xs[-1]
        left_src = max(0, x_min - 3)
        right_src = min(w - 1, x_max + 3)
        fill_color = (
            (bgr[y, left_src].astype(np.int32) + bgr[y, right_src].astype(np.int32))
            // 2
        ).astype(np.uint8)
        result[y, x_min : x_max + 1] = fill_color
    return result


def make_clean_background(bgr: np.ndarray, lines) -> np.ndarray:
    """
    使用改进的 inpaint 生成“无字版底图”：
    - 自适应主文字 mask
    - 扩展阴影/发光区域 mask 只用于 inpaint
    - 简单背景直接 inpaint，复杂背景先邻域填充再小半径 inpaint
    """
    if not lines:
        return bgr

    # 使用自适应或简单mask（主文字区域）
    if USE_ADAPTIVE_MASK:
        main_mask = build_adaptive_mask(bgr, lines)
    else:
        main_mask = build_text_mask_from_lines(bgr, lines)

    # 扩展阴影/发光区域，inpaint 时用这个大 mask
    shadow_mask = cv2.dilate(main_mask, np.ones((7, 7), np.uint8), iterations=2)

    is_simple = is_simple_background_region(bgr, shadow_mask)

    if is_simple:
        # 简单背景：直接 inpaint + 轻微模糊
        clean = cv2.inpaint(bgr, shadow_mask, INPAINT_RADIUS, cv2.INPAINT_TELEA)
    else:
        # 复杂背景：先用邻域像素粗填，再用小半径 NS 微调
        prefilled = fill_with_neighbor(bgr, shadow_mask)
        clean = cv2.inpaint(
            prefilled, shadow_mask, max(3, INPAINT_RADIUS // 2), cv2.INPAINT_NS
        )

    clean = cv2.GaussianBlur(clean, (3, 3), 0.5)

    # 只在 shadow_mask 区域应用 inpaint 结果
    result = bgr.copy()
    mask_3ch = cv2.cvtColor(shadow_mask, cv2.COLOR_GRAY2BGR) / 255.0
    result = (clean * mask_3ch + bgr * (1 - mask_3ch)).astype(np.uint8)

    return result


def ocr_images_to_ppt(
    image_paths: Sequence[str],
    output_pptx: str,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    use_text_color: bool = EXTRACT_TEXT_COLOR,
) -> str:
    """
    将图片通过OCR转换为可编辑文字的PPT（优化版）

    注意：该函数为内部实现，推荐通过 images_to_pdf_and_ppt /
    convert_images_dir_to_pdf_and_ppt 间接调用。
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    for idx, img_path in enumerate(image_paths, start=1):
        log.info(f"Processing slide #{idx}: {os.path.basename(img_path)}")

        bgr = read_bgr(img_path)

        # 预处理：放大和锐化
        ocr_img, scale = preprocess_for_ocr(bgr)

        if idx <= DEBUG_DUMP_FIRST_N:
            debug_dump(bgr, f"before_ocr_raw_{idx}")
            debug_dump(ocr_img, f"before_ocr_up_{idx}")
            log.info(f"slide#{idx} upscale scale={scale:.3f}")

        h0, w0 = bgr.shape[:2]  # 原图尺寸
        h1, w1 = ocr_img.shape[:2]  # OCR输入尺寸

        # 创建幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # OCR识别（PaddleOCR 直接处理 BGR 图像）
        lines = paddle_ocr(ocr_img)

        # 把bbox从OCR图坐标缩回原图坐标
        if lines and (w1 != w0 or h1 != h0):
            sx = w0 / float(w1)
            sy = h0 / float(h1)
            lines = [
                ([b[0] * sx, b[1] * sy, b[2] * sx, b[3] * sy], t, c)
                for (b, t, c) in lines
            ]

        # 再合并一轮
        y_tol = max(12, int(h0 * 0.008))
        x_gap = max(18, int(w0 * 0.01))
        lines = merge_lines(lines, y_tol=y_tol, x_gap=x_gap)

        # 统计正文行高，用于后续字号估计
        body_h_px = analyze_line_heights(lines)

        if not lines:
            log.warning(f"slide#{idx} no text detected")
        else:
            log.info(f"slide#{idx} detected {len(lines)} text boxes")

        # 估计背景颜色（用于颜色提取）
        bg_color = None
        if use_text_color and lines:
            bg_color = estimate_background_color(bgr, lines)
            if bg_color:
                log.info(f"slide#{idx} estimated background color: RGB{bg_color}")

        # 底图处理：可选 inpaint 生成"干净底图"
        bg_for_slide = bgr
        if add_background_image:
            if clean_background and lines:
                log.info(f"slide#{idx} applying inpainting...")
                bg_for_slide = make_clean_background(bgr, lines)
                if idx <= DEBUG_DUMP_FIRST_N:
                    debug_dump(bg_for_slide, f"clean_bg_{idx}")
            tmp = f"__ppt_bg_{idx}.png"
            add_background(slide, bg_for_slide, slide_w_emu, slide_h_emu, tmp)

        scale_x = slide_w_emu / w0
        scale_y = slide_h_emu / h0

        for bbox, text, conf in lines:
            x1, y1, x2, y2 = bbox
            if (x2 - x1) < 6 or (y2 - y1) < 6:
                continue

            # 计算字号
            font_size = estimate_font_pt(bbox, img_h_px=h0, body_h_px=body_h_px)

            # 文本框尺寸：直接使用OCR检测到的bbox尺寸
            bbox_width_emu = px_to_emu((x2 - x1), scale_x)
            bbox_height_emu = px_to_emu((y2 - y1), scale_y)

            width = bbox_width_emu
            height = bbox_height_emu

            left = px_to_emu(x1, scale_x)
            top = px_to_emu(y1, scale_y)

            # 添加透明文本框
            tb = slide.shapes.add_textbox(left, top, int(width), int(height))
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = False  # 禁用自动换行

            # 垂直居中（可选）
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 设置文本框透明
            tb.fill.background()  # 无填充
            tb.line.fill.background()  # 无边框

            p = tf.paragraphs[0]
            p.text = text

            # 使用原生分散对齐：多字符分散，单字符居中
            if len(text) > 1:
                p.alignment = PP_ALIGN.DISTRIBUTE
            else:
                p.alignment = PP_ALIGN.CENTER

            # 不再依赖 font.spacing 模拟分散，统一设为 0
            p.font.size = font_size
            p.font.spacing = Pt(0)

            # 提取并设置文字颜色
            if use_text_color:
                text_color = extract_text_color(bgr, bbox, bg_color)
                p.font.color.rgb = RGBColor(*text_color)
            else:
                # 默认黑色
                p.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(output_pptx)
    return output_pptx


# ----------------------------
# Public API
# ----------------------------


def images_to_pdf_and_ppt(
    image_paths: Sequence[str],
    output_pdf_path: Optional[str] = None,
    output_pptx_path: Optional[str] = None,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    extract_text_color: bool = EXTRACT_TEXT_COLOR,
) -> Dict[str, Optional[str]]:
    """
    将给定的一组图片转换为 PDF 和可编辑 PPTX。

    参数:
        image_paths: 按页面顺序排列的图片路径列表。
        output_pdf_path: 输出 PDF 文件路径，若为 None 则不生成 PDF。
        output_pptx_path: 输出 PPTX 文件路径，若为 None 则不生成 PPT。
        add_background_image: 是否在 PPT 中加入整页背景图。
        clean_background: 是否对背景进行 inpaint 处理（在 add_background_image 为 True 时生效）。
        extract_text_color: 是否根据原图估计文字颜色，用于 PPT 文本着色。

    返回:
        包含已生成文件路径的字典，例如:
        {
            "pdf": "/path/to/output.pdf" 或 None,
            "pptx": "/path/to/output_editable.pptx" 或 None,
        }
    """
    result: Dict[str, Optional[str]] = {"pdf": None, "pptx": None}

    if output_pdf_path is not None:
        result["pdf"] = images_to_pdf(image_paths, output_pdf_path)

    if output_pptx_path is not None:
        result["pptx"] = ocr_images_to_ppt(
            image_paths=image_paths,
            output_pptx=output_pptx_path,
            add_background_image=add_background_image,
            clean_background=clean_background,
            use_text_color=extract_text_color,
        )

    return result


def convert_images_dir_to_pdf_and_ppt(
    input_dir: str,
    output_pdf_path: Optional[str] = None,
    output_pptx_path: Optional[str] = None,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    extract_text_color: bool = EXTRACT_TEXT_COLOR,
) -> Dict[str, Optional[str]]:
    """
    给定图片目录，自动读取所有图片并生成 PDF + PPTX。

    参数:
        input_dir: 包含图片的目录，内部按文件名自然排序。
        其余参数同 images_to_pdf_and_ppt。

    返回:
        同 images_to_pdf_and_ppt。
    """
    image_paths = list_images_in_dir(input_dir)
    if not image_paths:
        raise ValueError(f"No images found in {input_dir!r}")

    return images_to_pdf_and_ppt(
        image_paths=image_paths,
        output_pdf_path=output_pdf_path,
        output_pptx_path=output_pptx_path,
        add_background_image=add_background_image,
        clean_background=clean_background,
        extract_text_color=extract_text_color,
    )


async def convert_images_dir_to_pdf_and_ppt_api(
    input_dir: str,
    output_pdf_path: Optional[str] = None,
    output_pptx_path: Optional[str] = None,
    api_url: Optional[str] = None,
    api_key: Optional[str] = None,
    model: Optional[str] = None,
    use_api_inpaint: bool = True,
    add_background_image: bool = ADD_BACKGROUND_IMAGE,
    clean_background: bool = CLEAN_BACKGROUND,
    use_text_color: bool = EXTRACT_TEXT_COLOR,
) -> Dict[str, Optional[str]]:
    """
    带 API inpainting 支持的图片转 PDF/PPTX 函数（异步版本）
    
    与 convert_images_dir_to_pdf_and_ppt 的区别：
    - 支持使用图像编辑 API 进行 inpainting（优先）
    - API 失败时自动 fallback 到传统 OpenCV inpaint
    - 支持重试机制（最多3次）
    
    参数:
        input_dir: 包含图片的目录，内部按文件名自然排序
        output_pdf_path: 输出 PDF 文件路径，若为 None 则不生成 PDF
        output_pptx_path: 输出 PPTX 文件路径，若为 None 则不生成 PPT
        api_url: 图像编辑 API 的 URL
        api_key: API 密钥
        model: 使用的模型名称
        use_api_inpaint: 是否启用 API inpainting（默认 True）
        add_background_image: 是否在 PPT 中加入整页背景图
        clean_background: 是否对背景进行 inpaint 处理
        extract_text_color: 是否根据原图估计文字颜色
    
    返回:
        包含已生成文件路径的字典
    """
    import asyncio
    from dataflow_agent.toolkits.imtool.req_img import generate_or_edit_and_save_image_async
    
    image_paths = list_images_in_dir(input_dir)
    if not image_paths:
        raise ValueError(f"No images found in {input_dir!r}")
    
    result: Dict[str, Optional[str]] = {"pdf": None, "pptx": None}
    
    # 生成 PDF
    if output_pdf_path is not None:
        result["pdf"] = images_to_pdf(image_paths, output_pdf_path)
    
    # 生成 PPTX（带 API inpainting 支持）
    if output_pptx_path is not None:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
        
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height
        
        for idx, img_path in enumerate(image_paths, start=1):
            log.info(f"Processing slide #{idx}: {os.path.basename(img_path)}")
            
            bgr = read_bgr(img_path)
            ocr_img, scale = preprocess_for_ocr(bgr)
            
            h0, w0 = bgr.shape[:2]
            h1, w1 = ocr_img.shape[:2]
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # OCR 识别
            lines = paddle_ocr(ocr_img)
            
            # 坐标映射
            if lines and (w1 != w0 or h1 != h0):
                sx = w0 / float(w1)
                sy = h0 / float(h1)
                lines = [
                    ([b[0] * sx, b[1] * sy, b[2] * sx, b[3] * sy], t, c)
                    for (b, t, c) in lines
                ]
            
            # 合并行
            y_tol = max(12, int(h0 * 0.008))
            x_gap = max(18, int(w0 * 0.01))
            lines = merge_lines(lines, y_tol=y_tol, x_gap=x_gap)
            
            body_h_px = analyze_line_heights(lines)
            bg_color = estimate_background_color(bgr, lines) if use_text_color and lines else None
            
            # 底图处理：优先使用 API inpainting
            bg_for_slide = bgr
            if add_background_image:
                if clean_background and lines and use_api_inpaint and api_url and api_key and model:
                    # 使用 API inpainting（带重试）
                    async def _call_inpaint_api_with_retry(retries: int = 3, delay: float = 1.0) -> bool:
                        last_err: Optional[Exception] = None
                        for attempt in range(1, retries + 1):
                            try:
                                await generate_or_edit_and_save_image_async(
                                    prompt=inpaint_prompt,
                                    save_path=clean_bg_path,
                                    aspect_ratio="16:9",
                                    api_url=api_url,
                                    api_key=api_key,
                                    model=model,
                                    image_path=temp_img_path,
                                    use_edit=True,
                                )
                                return True
                            except Exception as e:
                                last_err = e
                                log.error(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} inpainting attempt {attempt}/{retries} failed: {e}")
                                if attempt < retries:
                                    try:
                                        await asyncio.sleep(delay)
                                    except Exception:
                                        pass
                        log.error(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} inpainting failed after {retries} attempts: {last_err}")
                        return False
                    
                    try:
                        # 生成文字掩码
                        text_mask = build_adaptive_mask(bgr, lines)
                        
                        # 保存临时文件
                        import tempfile
                        with tempfile.TemporaryDirectory() as tmpdir:
                            temp_img_path = os.path.join(tmpdir, f"temp_{idx}.png")
                            clean_bg_path = os.path.join(tmpdir, f"clean_{idx}.png")
                            cv2.imwrite(temp_img_path, bgr)
                            
                            # 构造 inpainting 提示词
                            inpaint_prompt = "请智能修复图像中文字被移除后的区域，保持背景的连续性、一致性和自然过渡，使修复后的图像看起来完整无缺，并且原图涉及的图标你需要尽量保留；"
                            
                            log.info(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} 开始调用图像编辑API进行inpainting（最多重试3次）...")
                            
                            api_success = await _call_inpaint_api_with_retry(retries=3, delay=1.0)
                            
                            if api_success and os.path.exists(clean_bg_path):
                                bg_for_slide = read_bgr(clean_bg_path)
                                log.info(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} API inpainting成功")
                            else:
                                log.warning(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} API inpainting失败，使用本地inpaint")
                                bg_for_slide = make_clean_background(bgr, lines)
                    except Exception as e:
                        log.error(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} inpainting流程失败: {e}，使用本地inpaint")
                        try:
                            bg_for_slide = make_clean_background(bgr, lines)
                        except Exception as e2:
                            log.error(f"[convert_images_dir_to_pdf_and_ppt_api] slide#{idx} 本地inpaint也失败: {e2}，使用原图")
                            bg_for_slide = bgr
                elif clean_background and lines:
                    # 不使用 API，直接使用本地 inpaint
                    log.info(f"slide#{idx} applying local inpainting...")
                    bg_for_slide = make_clean_background(bgr, lines)
                
                tmp = f"__ppt_bg_{idx}.png"
                add_background(slide, bg_for_slide, slide_w_emu, slide_h_emu, tmp)
            
            # 添加文本框（与原函数相同的逻辑）
            scale_x = slide_w_emu / w0
            scale_y = slide_h_emu / h0
            
            for bbox, text, conf in lines:
                x1, y1, x2, y2 = bbox
                if (x2 - x1) < 6 or (y2 - y1) < 6:
                    continue

                font_size = estimate_font_pt(bbox, img_h_px=h0, body_h_px=body_h_px)

                bbox_width_emu = px_to_emu((x2 - x1), scale_x)
                bbox_height_emu = px_to_emu((y2 - y1), scale_y)
                width = bbox_width_emu
                height = bbox_height_emu
                left = px_to_emu(x1, scale_x)
                top = px_to_emu(y1, scale_y)

                tb = slide.shapes.add_textbox(left, top, int(width), int(height))
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = False
                tb.fill.background()
                tb.line.fill.background()

                # 垂直居中（可选）
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                p = tf.paragraphs[0]
                p.text = text

                # 使用原生分散对齐
                if len(text) > 1:
                    p.alignment = PP_ALIGN.DISTRIBUTE
                else:
                    p.alignment = PP_ALIGN.CENTER

                p.font.size = font_size
                p.font.spacing = Pt(0)

                if use_text_color:
                    text_color = extract_text_color(bgr, bbox, bg_color)
                    p.font.color.rgb = RGBColor(*text_color)
                else:
                    p.font.color.rgb = RGBColor(0, 0, 0)
        
        prs.save(output_pptx_path)
        result["pptx"] = output_pptx_path
    
    return result


if __name__ == "__main__":
    """
    简单本地测试入口：
    - 直接运行本文件即可测试 PaddleOCR 对指定图片的识别效果
    - 识别结果会打印在终端，并把画好检测框的图片保存到指定路径
    """
    # 测试图片路径（也是可视化输出路径）
    img_path = f"{get_project_root()}/tests/test_02.png"

    if not os.path.exists(img_path):
        raise FileNotFoundError(f"测试图片不存在: {img_path}")

    # 调用封装好的单页接口
    info = paddle_ocr_page_with_layout(img_path)

    print("=== PaddleOCR 测试结果 ===")
    print(f"image_size: {info['image_size']}")
    print(f"body_h_px: {info['body_h_px']}")
    print(f"bg_color: {info['bg_color']}")
    print(f"检测到文本框数量: {len(info['lines'])}")

    for i, (bbox, text, conf) in enumerate(info["lines"], start=1):
        print(f"[{i:02d}] conf={conf:.1f} bbox={bbox} text={text}")

    # 把检测框画在图上，并保存到文件，而不是弹出窗口
    try:
        bgr = read_bgr(img_path)
        vis = bgr.copy()
        for bbox, text, conf in info["lines"]:
            x1, y1, x2, y2 = map(int, bbox)
            cv2.rectangle(vis, (x1, y1), (x2, y2), (0, 255, 0), 2)

        save_path = f"{get_project_root()}/tests/test_01_paddle_frame.png"
        ok = cv2.imwrite(save_path, vis)
        if ok:
            log.info(f"PaddleOCR 可视化结果已保存到: {save_path}")
        else:
            log.warning(f"PaddleOCR 可视化结果保存失败: {save_path}")
    except Exception as e:
        log.warning(f"可视化失败: {e}")
        # 不影响纯文本打印结果
