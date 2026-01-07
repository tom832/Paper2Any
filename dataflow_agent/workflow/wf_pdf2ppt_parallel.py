"""
pdf2ppt_with_sam workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
基于 slides PDF:
1. 将 PDF 每页渲染为 PNG
2. 对每页图片用 PaddleOCR 做文字 OCR
3. 对每页图片用 MinerU 做版面分析（区分 Text vs Image/Table）
4. 对每页图片用 SAM 做图标 / 图块分割
5. 智能合并：
   - MinerU 划定 "图表区" (Image/Table) 和 "正文区"。
   - OCR 文本如果落在 "图表区" 则丢弃，防止图片上的文字重复生成。
   - SAM 图块如果落在 "图表区" 则丢弃（由 MinerU 负责）；如果在 "正文区" 且包含文字则丢弃（防止把文字当图）；
     剩下的 SAM 块被视为 "无字图标"，进行抠图后保留。
   - MinerU 提取的图片直接复用其 sub_images 目录，不再手动裁剪。
   - 字体归一化：全局统计正文和标题字号，强制统一，保证整齐。
   - 使用 AI Inpainting 生成干净背景。
"""

from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional
from collections import Counter

import cv2
import numpy as np
import fitz  # PyMuPDF
import yaml
from PIL import Image

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.state import Paper2FigureState 
from dataflow_agent.utils import get_project_root

# Tools
from dataflow_agent.toolkits.imtool.sam_tool import segment_layout_boxes, segment_layout_boxes_server, free_sam_model
from dataflow_agent.toolkits.imtool.bg_tool import local_tool_for_bg_remove, free_bg_rm_model
from dataflow_agent.toolkits.imtool.mineru_tool import recursive_mineru_layout
from dataflow_agent.toolkits.imtool.req_img import gemini_multi_image_edit_async
from dataflow_agent.toolkits.imtool import ppt_tool

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

log = get_logger(__name__)

# Load configuration from yaml
def load_server_config():
    root = get_project_root()
    config_path = root / "conf" / "model_servers.yaml"
    if not config_path.exists():
        log.warning(f"Config file not found at {config_path}, using defaults.")
        return {}
    try:
        with open(config_path, "r") as f:
            return yaml.safe_load(f) or {}
    except Exception as e:
        log.error(f"Failed to load config: {e}")
        return {}

SERVER_CONFIG = load_server_config()

# Helper to construct URLs
def get_sam_urls():
    # Check env var first
    if os.environ.get("SAM_SERVER_URLS"):
        return os.environ.get("SAM_SERVER_URLS").split(",")
    
    # Try config
    sam_cfg = SERVER_CONFIG.get("sam", {})
    instances = sam_cfg.get("instances", [])
    if instances:
        urls = []
        for inst in instances:
            for port in inst.get("ports", []):
                urls.append(f"http://127.0.0.1:{port}")
        if urls:
            return urls
            
    # Default
    return ["http://localhost:8021", "http://localhost:8022","http://localhost:8023"]

def get_ocr_urls():
    # Check env var first
    if os.environ.get("OCR_SERVER_URLS"):
        return os.environ.get("OCR_SERVER_URLS").split(",")

    # Try config
    ocr_cfg = SERVER_CONFIG.get("ocr", {})
    if ocr_cfg:
        host = ocr_cfg.get("host", "0.0.0.0")
        if host == "0.0.0.0": host = "127.0.0.1"
        port = ocr_cfg.get("port", 8003)
        return [f"http://{host}:{port}"]

    # Default
    return ["http://localhost:8003"]

SAM_SERVER_URLS = get_sam_urls()
OCR_SERVER_URLS = get_ocr_urls()


def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    为本次 pdf2ppt_with_sam workflow 创建统一的输出目录：
    - 如果 state.result_path 已存在，直接使用；
    - 否则使用项目根目录下 outputs/pdf2ppt_with_sam/<timestamp>。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(__import__("time").time())
    base_dir = (root / "outputs" / "pdf2ppt_with_sam" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path


def _run_sam_on_pages(image_paths: List[str], base_dir: str) -> List[Dict[str, Any]]:
    """
    对每一页图片运行 SAM，输出 layout_items。
    """
    results: List[Dict[str, Any]] = []
    sam_ckpt = f"{get_project_root()}/sam_b.pt"

    for page_idx, img_path in enumerate(image_paths):
        img_path_obj = Path(img_path)
        if not img_path_obj.exists():
            log.warning(f"[pdf2ppt_with_sam] image not found for SAM: {img_path}")
            results.append({"page_idx": page_idx, "layout_items": []})
            continue

        out_dir = Path(base_dir) / "layout_items" / f"page_{page_idx+1:03d}"
        out_dir.mkdir(parents=True, exist_ok=True)

        # 1. SAM 分割 (使用远程服务)
        try:
            layout_items = segment_layout_boxes_server(
                image_path=str(img_path_obj),
                output_dir=str(out_dir),
                server_urls=SAM_SERVER_URLS,
                checkpoint=sam_ckpt,
                min_area=200,
                min_score=0.0,
                iou_threshold=0.4,
                top_k=25,
                nms_by="mask",
            )
        except Exception as e:
            log.error(f"[pdf2ppt_with_sam] Remote SAM failed: {e}. Fallback to local.")
            # Fallback to local if server fails
            layout_items = segment_layout_boxes(
                image_path=str(img_path_obj),
                output_dir=str(out_dir),
                checkpoint=sam_ckpt,
                min_area=200,
                min_score=0.0,
                iou_threshold=0.4,
                top_k=25,
                nms_by="mask",
            )
            
        log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] SAM found {len(layout_items)} items")

        # 2. 映射 bbox 到像素坐标（基于整页尺寸）
        try:
            pil_img = Image.open(str(img_path_obj))
            w, h = pil_img.size
        except Exception as e:
            log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] open image failed: {e}")
            w, h = 1024, 768

        for it in layout_items:
            bbox = it.get("bbox")
            if bbox and len(bbox) == 4:
                x1n, y1n, x2n, y2n = bbox
                x1 = int(round(x1n * w))
                y1 = int(round(y1n * h))
                x2 = int(round(x2n * w))
                y2 = int(round(y2n * h))
                if x2 > x1 and y2 > y1:
                    it["bbox_px"] = [x1, y1, x2, y2]

        results.append({"page_idx": page_idx, "layout_items": layout_items})

    # 显式释放 SAM 模型
    try:
        free_sam_model(checkpoint=sam_ckpt)
    except Exception as e:
        log.error(f"[pdf2ppt_with_sam] free_sam_model failed: {e}")

    return results


@register("pdf2ppt_parallel")
def create_pdf2ppt_with_sam_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf pdf2ppt_with_sam
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState, entry_point="_start_")

    # ==============================
    # NODES
    # ==============================

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        _ensure_result_path(state)
        return state

    async def pdf_to_images_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        将 PDF 每一页渲染为 PNG。
        """
        pdf_path = getattr(state, "pdf_file", None)
        if not pdf_path:
            log.error("[pdf2ppt_with_sam] state.pdf_file is empty")
            return state

        base_dir = Path(_ensure_result_path(state))
        img_dir = base_dir / "slides_png"
        image_paths = ppt_tool.pdf_to_images(pdf_path, str(img_dir))
        state.slide_images = image_paths
        return state

    async def slides_ocr_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        对每一页图片用 PaddleOCR 做 OCR。
        使用 asyncio.to_thread 包装同步调用，避免阻塞事件循环。
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.error("[pdf2ppt_with_sam] no slide_images for OCR")
            return state

        def _sync_ocr_all_pages():
            """同步执行所有页面的 OCR"""
            ocr_pages: List[Dict[str, Any]] = []
            for page_idx, img_path in enumerate(image_paths):
                try:
                    # 优先使用远程 OCR 服务
                    try:
                        result = ppt_tool.paddle_ocr_page_with_layout_server(img_path, server_urls=OCR_SERVER_URLS)
                    except Exception as e:
                        log.warning(f"[pdf2ppt_with_sam][OCR] remote failed: {e}. Fallback to local.")
                        result = ppt_tool.paddle_ocr_page_with_layout(img_path)
                except Exception as e:
                    log.error(f"[pdf2ppt_with_sam][OCR] page#{page_idx+1} failed: {e}")
                    result = {
                        "image_size": None,
                        "lines": [],
                        "body_h_px": None,
                        "bg_color": None,
                        "path": img_path,
                        "page_idx": page_idx,
                    }
                result["page_idx"] = page_idx
                result["path"] = img_path
                ocr_pages.append(result)
            return ocr_pages

        # 在线程池中执行同步 OCR，不阻塞事件循环
        ocr_pages = await asyncio.to_thread(_sync_ocr_all_pages)
        state.ocr_pages = ocr_pages
        return state

    async def slides_mineru_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        对每一页 PNG 使用 MinerU 做版面识别：
        - 输出每页的 mineru_items，包含 type / bbox(norm) / text 等
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.error("[pdf2ppt_with_sam] no slide_images for MinerU")
            return state

        base_dir = Path(_ensure_result_path(state))
        mineru_dir = base_dir / "mineru_pages"
        mineru_dir.mkdir(parents=True, exist_ok=True)

        # MinerU 端口，优先从 state.request.mineru_port 读取
        # MinerU LB Port 8010
        port = getattr(getattr(state, "request", None), "mineru_port", 8010)
        # 复杂度深度可从 state 或常量
        max_depth = getattr(state, "mask_detail_level", 3)

        mineru_pages: List[Dict[str, Any]] = []

        for page_idx, img_path in enumerate(image_paths):
            try:
                out_dir = mineru_dir / f"page_{page_idx+1:03d}"
                out_dir.mkdir(parents=True, exist_ok=True)

                log.critical(f"【mineru node】:  {out_dir}")

                mineru_items = await recursive_mineru_layout(
                    image_path=str(img_path),
                    port=port,
                    max_depth=3,
                    output_dir=str(out_dir),
                )
                
                # 记录 MinerU 输出目录，方便后续找 sub_images
                # recursive_mineru_layout 会在 out_dir 下直接输出或创建子目录
                # 这里我们记录 out_dir，后续可以在里面找 sub_images
                
                mineru_pages.append({
                    "page_idx": page_idx,
                    "blocks": mineru_items,
                    "path": img_path,
                    "mineru_output_dir": str(out_dir)
                })
                log.info(f"[pdf2ppt_with_sam][MinerU] page#{page_idx+1} got {len(mineru_items)} blocks")
            except Exception as e:
                log.error(f"[pdf2ppt_with_sam][MinerU] page#{page_idx+1} failed: {e}")
                mineru_pages.append({
                    "page_idx": page_idx,
                    "blocks": [],
                    "path": img_path,
                })

        state.mineru_pages = mineru_pages

        log.critical(f"[state.mineru_pages]:  {state.mineru_pages}")

        return state

    async def slides_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        对每一页图片运行 SAM 用于图标 / 图块分割。
        使用 asyncio.to_thread 包装同步调用，避免阻塞事件循环。
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.error("[pdf2ppt_with_sam] no slide_images for SAM")
            return state

        base_dir = _ensure_result_path(state)
        
        # 在线程池中执行同步 SAM，不阻塞事件循环
        sam_pages = await asyncio.to_thread(_run_sam_on_pages, image_paths, base_dir)
        state.sam_pages = sam_pages
        return state

    async def slides_layout_bg_remove_node(state: Paper2FigureState, sam_pages: List[Dict[str, Any]] = None) -> Paper2FigureState:
        """
        对每一页 SAM layout PNG 做背景抠图：
        - 输入: state.sam_pages[*].layout_items[].png_path 或传入的 sam_pages
        - 输出: 为每个 layout_item 写入 fg_png_path（抠完背景的 PNG）
        使用 asyncio.to_thread 包装同步调用，避免阻塞事件循环。
        """
        # 支持从参数传入 sam_pages（用于并行分支）
        if sam_pages is None:
            sam_pages = getattr(state, "sam_pages", []) or []
        
        if not sam_pages:
            log.error("[pdf2ppt_with_sam] no sam_pages for bg remove")
            return state

        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "sam_icons"
        icons_dir.mkdir(parents=True, exist_ok=True)

        model_path = getattr(getattr(state, "request", None), "bg_rm_model", None)

        def _sync_bg_remove():
            """同步执行所有背景移除"""
            processed = 0
            for p in sam_pages:
                page_idx = p.get("page_idx", 0)
                for it in p.get("layout_items", []):
                    png_path = it.get("png_path")
                    if not png_path or not os.path.exists(png_path):
                        continue
                    
                    # 背景抠图 - 添加页码前缀避免文件名冲突
                    try:
                        # 从原始路径提取文件名
                        original_stem = Path(png_path).stem
                        # 创建带页码的输出文件名
                        output_filename = f"page_{page_idx+1:03d}_{original_stem}_bg_removed.png"
                        output_path = icons_dir / output_filename
                        
                        req = {
                            "image_path": png_path,
                            "output_dir": str(icons_dir),
                        }
                        if model_path:
                            req["model_path"] = model_path
                        
                        fg_path = local_tool_for_bg_remove(req)
                        
                        # 重命名文件以包含页码
                        if fg_path and os.path.exists(fg_path):
                            # 将生成的文件重命名为带页码的文件名
                            fg_path_obj = Path(fg_path)
                            if fg_path_obj.name != output_filename:
                                new_fg_path = fg_path_obj.parent / output_filename
                                fg_path_obj.rename(new_fg_path)
                                fg_path = str(new_fg_path)
                            
                            it["fg_png_path"] = fg_path
                        else:
                            it["fg_png_path"] = png_path
                        
                        processed += 1
                    except Exception as e:
                        log.error(f"[pdf2ppt_with_sam][bg_rm] failed for {png_path}: {e}")
                        it["fg_png_path"] = png_path

            # 抠图完成后可尝试释放模型（忽略失败）
            try:
                if model_path:
                    free_bg_rm_model(model_path=model_path)
            except Exception as e:
                log.error(f"[pdf2ppt_with_sam] free_bg_rm_model failed: {e}")
            
            return processed

        # 在线程池中执行同步背景移除，不阻塞事件循环
        processed = await asyncio.to_thread(_sync_bg_remove)

        log.info(f"[pdf2ppt_with_sam] bg remove processed: {processed} items")
        
        # 将处理后的 sam_pages 写回 state
        state.sam_pages = sam_pages
        return state

    # ==============================================================
    # 并行处理节点：同时执行 OCR、MinerU、SAM+背景移除
    # ==============================================================
    async def parallel_processing_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        并行执行三个分支：
        1. slides_ocr_node -> ocr_pages
        2. slides_mineru_node -> mineru_pages
        3. slides_sam_node + slides_layout_bg_remove_node -> sam_pages
        
        三个分支共享 state.slide_images 作为输入，各自写入不同的输出字段。
        """
        import copy
        import time
        
        log.info("[parallel_processing] 开始并行处理 OCR / MinerU / SAM+BgRemove ...")
        start_time = time.time()
        
        # 定义三个分支任务
        async def ocr_branch():
            """OCR 分支"""
            log.info("[parallel_processing][OCR] 分支启动")
            branch_state = copy.copy(state)  # 浅拷贝，共享 slide_images
            result = await slides_ocr_node(branch_state)
            log.info(f"[parallel_processing][OCR] 分支完成，提取了 {len(getattr(result, 'ocr_pages', []))} 页")
            return ("ocr", result)
        
        async def mineru_branch():
            """MinerU 分支"""
            log.info("[parallel_processing][MinerU] 分支启动")
            branch_state = copy.copy(state)
            result = await slides_mineru_node(branch_state)
            log.info(f"[parallel_processing][MinerU] 分支完成，提取了 {len(getattr(result, 'mineru_pages', []))} 页")
            return ("mineru", result)
        
        async def sam_branch():
            """SAM + 背景移除 分支（串行执行）"""
            log.info("[parallel_processing][SAM] 分支启动")
            branch_state = copy.copy(state)
            
            # 先执行 SAM
            branch_state = await slides_sam_node(branch_state)
            sam_pages = getattr(branch_state, "sam_pages", [])
            log.info(f"[parallel_processing][SAM] SAM 完成，提取了 {len(sam_pages)} 页")
            
            # 再执行背景移除
            branch_state = await slides_layout_bg_remove_node(branch_state, sam_pages=sam_pages)
            log.info("[parallel_processing][SAM] 背景移除完成")
            
            return ("sam", branch_state)
        
        # 并行执行三个分支
        results = await asyncio.gather(
            ocr_branch(),
            mineru_branch(),
            sam_branch(),
            return_exceptions=True
        )
        
        # 合并结果到 state
        for r in results:
            if isinstance(r, Exception):
                log.error(f"[parallel_processing] 分支执行失败: {r}")
                import traceback
                traceback.print_exc()
                continue
            
            branch_name, branch_state = r
            
            if branch_name == "ocr":
                ocr_pages = getattr(branch_state, "ocr_pages", None)
                if ocr_pages:
                    state.ocr_pages = ocr_pages
                    log.info(f"[parallel_processing] 合并 OCR 结果: {len(ocr_pages)} 页")
            
            elif branch_name == "mineru":
                mineru_pages = getattr(branch_state, "mineru_pages", None)
                if mineru_pages:
                    state.mineru_pages = mineru_pages
                    log.info(f"[parallel_processing] 合并 MinerU 结果: {len(mineru_pages)} 页")
            
            elif branch_name == "sam":
                sam_pages = getattr(branch_state, "sam_pages", None)
                if sam_pages:
                    state.sam_pages = sam_pages
                    log.info(f"[parallel_processing] 合并 SAM 结果: {len(sam_pages)} 页")
        
        elapsed = time.time() - start_time
        log.info(f"[parallel_processing] 并行处理完成，耗时 {elapsed:.2f}s")
        
        return state

    async def slides_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        结合 MinerU + OCR + SAM 结果生成可编辑 PPT：
        
        改进点：
        1. MinerU 图片渲染修复：优先复用 MinerU 输出目录下的 sub_images，无法匹配时再手动裁剪。
        2. 字体归一化：
           - 统计全页正文（Body）文本的平均字号，取众数作为标准正文字号。
           - 标题（Title）字号设为标准正文的 1.5 倍（或取 Title 众数）。
           - 强制所有 Body 文本使用 Standard Body Font，所有 Title 文本使用 Standard Title Font。
        3. 背景生成开关：
           - 使用 state.use_ai_edit 控制是否调用 AI 生成纯净背景；
           - 关闭时直接使用纯白背景。
        4. 并行 API 调用：
           - 将 Inpainting API 调用改为并行执行，加快多页处理速度。
        """
        
        ocr_pages: List[Dict[str, Any]] = getattr(state, "ocr_pages", []) or []
        sam_pages: List[Dict[str, Any]] = getattr(state, "sam_pages", []) or []
        mineru_pages: List[Dict[str, Any]] = getattr(state, "mineru_pages", []) or []

        if not ocr_pages:
            log.error("[pdf2ppt_with_sam] no ocr_pages, abort PPT generation")
            return state

        # 建立索引
        sam_dict = {p.get("page_idx", 0): p.get("layout_items", []) for p in sam_pages}
        
        # mineru_dict 存放 {"blocks": [], "mineru_output_dir": ...}
        # 修复：为了防止 page_idx 类型不一致 (int vs str)，构建更鲁棒的索引
        mineru_dict = {}
        for p in mineru_pages:
            pid = p.get("page_idx", 0)
            mineru_dict[pid] = p        # 原始类型
            mineru_dict[str(pid)] = p   # 字符串类型兼容

        # 以 PPT 工具里的默认比例创建 Presentation
        prs = Presentation()
        prs.slide_width = Inches(ppt_tool.SLIDE_W_IN)
        prs.slide_height = Inches(ppt_tool.SLIDE_H_IN)
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height

        # 初始化 base_dir，确保后续逻辑都能访问
        base_dir = Path(_ensure_result_path(state))

        # ==========================================================
        # 辅助函数：API 重试逻辑
        # ==========================================================
        async def _call_image_api_with_retry(coro_factory, retries: int = 3, delay: float = 1.0) -> bool:
            """
            对图像生成/编辑进行最多 retries 次重试。
            """
            last_err: Optional[Exception] = None
            for attempt in range(1, retries + 1):
                try:
                    await coro_factory()
                    return True
                except Exception as e:
                    last_err = e
                    log.error(f"[pdf2ppt_with_sam] image api failed attempt {attempt}/{retries}: {e}")
                    if attempt < retries:
                        try:
                            await asyncio.sleep(delay)
                        except Exception:
                            pass
            log.error(f"[pdf2ppt_with_sam] image api failed after {retries} attempts: {last_err}")
            return False

        # ==========================================================
        # 辅助函数：字体和几何计算
        # ==========================================================
        def _bbox_area(bbox):
            return max(0, bbox[2] - bbox[0]) * max(0, bbox[3] - bbox[1])

        def _get_intersection_area(bbox1, bbox2):
            x1 = max(bbox1[0], bbox2[0])
            y1 = max(bbox1[1], bbox2[1])
            x2 = min(bbox1[2], bbox2[2])
            y2 = min(bbox1[3], bbox2[3])
            return max(0, x2 - x1) * max(0, y2 - y1)

        def _is_inside(inner, outer, threshold=0.9):
            inter = _get_intersection_area(inner, outer)
            inner_a = _bbox_area(inner)
            if inner_a <= 0: return False
            return (inter / inner_a) >= threshold

        def _is_overlap(bbox1, bbox2, threshold=0.1):
            inter = _get_intersection_area(bbox1, bbox2)
            min_area = min(_bbox_area(bbox1), _bbox_area(bbox2))
            if min_area <= 0: return False
            return (inter / min_area) >= threshold

        # ==========================================================
        # Phase 1: 准备渲染数据 & 创建 AI 任务
        # ==========================================================
        
        pages_render_data: List[Dict[str, Any]] = []
        ai_coroutines = []  # List of awaitables
        
        # 循环处理每一页的布局分析
        for pinfo in ocr_pages:
            page_idx = pinfo.get("page_idx", 0)
            
            # 兼容性查找
            mineru_page_data = mineru_dict.get(page_idx)
            if not mineru_page_data:
                mineru_page_data = mineru_dict.get(str(page_idx), {})
                if mineru_page_data:
                    log.warning(f"[pdf2ppt_with_sam] page_idx mismatch fixed by str conversion: {page_idx}")
            
            img_path = pinfo.get("path")
            lines = pinfo.get("lines", []) # List of (bbox, text, conf)
            
            if not img_path or not os.path.exists(img_path):
                log.warning(f"[pdf2ppt_with_sam] missing img for page#{page_idx+1}: {img_path}")
                continue

            # 读取原始图像信息
            try:
                pil_img = Image.open(img_path)
                w0, h0 = pil_img.size
            except Exception as e:
                log.error(f"Failed to open image {img_path}: {e}")
                continue

            # -----------------------------------------------------------
            # Step 1: 分析 MinerU 结果，划定 "Image Zone" 并找回 sub_images
            # -----------------------------------------------------------
            mineru_blocks = mineru_page_data.get("blocks", [])
            mineru_out_dir = mineru_page_data.get("mineru_output_dir")
            
            image_zones = []  # List of {"bbox": [x1,y1,x2,y2], "type": str, "img_path": str}
            
            # 尝试定位 sub_images 目录
            sub_images_dir = None
            sub_images_dirs: List[Path] = []
            if mineru_out_dir:
                try:
                    page_root = Path(mineru_out_dir)
                    direct = page_root / "sub_images"
                    if direct.exists() and direct.is_dir():
                        sub_images_dirs.append(direct)
                    for d in page_root.rglob("sub_images"):
                        if d.is_dir():
                            sub_images_dirs.append(d)
                    seen = set()
                    unique_dirs: List[Path] = []
                    for d in sub_images_dirs:
                        rp = str(d.resolve())
                        if rp not in seen:
                            seen.add(rp)
                            unique_dirs.append(d)
                    for d in unique_dirs:
                        pngs = list(d.glob("*.png"))
                        if pngs:
                            sub_images_dir = d
                            break
                    if sub_images_dir:
                        sub_files = sorted([p.name for p in sub_images_dir.glob("*.png")])
                        log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] MinerU sub_images dir: {sub_images_dir}, found {len(sub_files)} pngs")
                except Exception as e:
                    log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] search sub_images failed: {e}")

            for idx, blk in enumerate(mineru_blocks):
                btype = (blk.get("type") or "").lower()
                bbox = blk.get("bbox")  # norm
                if not bbox or len(bbox) != 4:
                    continue
                
                x1 = int(round(bbox[0] * w0))
                y1 = int(round(bbox[1] * h0))
                x2 = int(round(bbox[2] * w0))
                y2 = int(round(bbox[3] * h0))
                
                if x2 <= x1 or y2 <= y1: continue
                px_bbox = [x1, y1, x2, y2]
                
                is_image_zone = btype in ['image', 'figure', 'table', 'formula']
                img_path_found = None
                
                if is_image_zone:
                    if blk.get("img_path") and os.path.exists(blk["img_path"]):
                        img_path_found = blk["img_path"]
                    
                    if not img_path_found and sub_images_dir:
                        try:
                            depth = blk.get("depth", 0)
                            try:
                                depth = int(depth)
                            except Exception:
                                depth = 0
                            prefix = f"depth{depth}_blk{idx}_"
                            for f in sorted(sub_images_dir.glob("*.png")):
                                if f.name.startswith(prefix):
                                    img_path_found = str(f.resolve())
                                    break
                        except Exception as e:
                            log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] match sub_images failed: {e}")

                    if not img_path_found:
                        fallback_dir = base_dir / "mineru_fallback_crops" / f"page_{page_idx+1:03d}"
                        fallback_dir.mkdir(parents=True, exist_ok=True)
                        save_path = fallback_dir / f"mineru_{idx}_{btype}.png"
                        try:
                            if not save_path.exists():
                                crop = pil_img.crop((x1, y1, x2, y2))
                                crop.save(save_path)
                            img_path_found = str(save_path)
                        except Exception as e:
                            log.error(f"Failed to crop mineru block {idx}: {e}")

                    if img_path_found:
                        image_zones.append({
                            "bbox": px_bbox,
                            "type": btype,
                            "img_path": img_path_found
                        })

            # -----------------------------------------------------------
            # Step 2: 过滤 OCR 文字
            # -----------------------------------------------------------
            final_ocr_lines = [] # (bbox, text, conf, type, raw_pt)
            
            for line in lines:
                l_bbox, l_text, l_conf = line
                is_in_image = False
                for zone in image_zones:
                    if _is_inside(l_bbox, zone["bbox"]):
                        is_in_image = True
                        break
                
                if not is_in_image:
                    l_type = "body"
                    for blk in mineru_blocks:
                        btype = (blk.get("type") or "").lower()
                        b_bbox = blk.get("bbox")
                        if not b_bbox: continue
                        bx1 = int(round(b_bbox[0] * w0))
                        by1 = int(round(b_bbox[1] * h0))
                        bx2 = int(round(b_bbox[2] * w0))
                        by2 = int(round(b_bbox[3] * h0))
                        
                        if btype in ['title', 'header'] and _is_inside(l_bbox, [bx1, by1, bx2, by2]):
                            l_type = "title"
                            break
                    
                    # 预先计算原始字号，方便后续聚类
                    raw_pt_obj = ppt_tool.estimate_font_pt(l_bbox, img_h_px=h0, body_h_px=None)
                    raw_pt = raw_pt_obj.pt if hasattr(raw_pt_obj, "pt") else raw_pt_obj

                    final_ocr_lines.append((l_bbox, l_text, l_conf, l_type, raw_pt))

            # -----------------------------------------------------------
            # Step 3: 过滤 SAM 图块
            # -----------------------------------------------------------
            raw_sam_items = sam_dict.get(page_idx, [])
            final_sam_items = []
            
            for item in raw_sam_items:
                s_bbox = item.get("bbox_px")
                if not s_bbox: continue
                is_in_image = False
                for zone in image_zones:
                    if _is_inside(s_bbox, zone["bbox"], threshold=0.6):
                        is_in_image = True
                        break
                if is_in_image: continue

                is_text_block = False
                for line in final_ocr_lines:
                    l_bbox = line[0]
                    if _is_overlap(s_bbox, l_bbox, threshold=0.3) or _is_inside(l_bbox, s_bbox):
                        is_text_block = True
                        break
                if is_text_block: continue

                w = s_bbox[2] - s_bbox[0]
                h = s_bbox[3] - s_bbox[1]
                if w < 5 or h < 5: continue
                if w*h < 400: continue 
                
                final_sam_items.append(item)

            # -----------------------------------------------------------
            # Step 4: 准备 AI 背景生成任务
            # -----------------------------------------------------------
            clean_bg_path = base_dir / "clean_backgrounds" / f"clean_bg_{page_idx+1:03d}.png"
            clean_bg_path.parent.mkdir(parents=True, exist_ok=True)
            
            use_ai_bg = bool(getattr(state, "use_ai_edit", False))
            log.critical(f"[pdf2ppt 是否使用AI： ][page#{page_idx+1}] use_ai_bg={use_ai_bg}")
            
            ai_task = None
            if use_ai_bg and os.path.exists(img_path):
                try:
                    # A. 生成 Mask (黑底白框)
                    ori_cv = cv2.imread(img_path)
                    if ori_cv is not None:
                        h_cv, w_cv = ori_cv.shape[:2]
                        mask_cv = np.zeros((h_cv, w_cv), dtype=np.uint8)  # 黑底
                        
                        # 绘制 OCR 区域 (白框)
                        for line in final_ocr_lines:
                            bbox = line[0]
                            pad = 5
                            mx1 = int(max(0, bbox[0] - pad))
                            my1 = int(max(0, bbox[1] - pad))
                            mx2 = int(min(w_cv, bbox[2] + pad))
                            my2 = int(min(h_cv, bbox[3] + pad))
                            cv2.rectangle(mask_cv, (mx1, my1), (mx2, my2), (255), -1)
                        
                        mask_path = base_dir / "masks" / f"mask_{page_idx+1:03d}.png"
                        mask_path.parent.mkdir(parents=True, exist_ok=True)
                        cv2.imwrite(str(mask_path), mask_cv)
                        
                        # B. 准备 AI 调用闭包
                        req_cfg = getattr(state, "request", None) or {}
                        if not isinstance(req_cfg, dict):
                            req_cfg = req_cfg.__dict__ if hasattr(req_cfg, "__dict__") else {}
                            
                        api_key = req_cfg.get("api_key") or os.getenv("DF_API_KEY")
                        api_url = req_cfg.get("chat_api_url") or "https://api.apiyi.com"
                        model_name = req_cfg.get("gen_fig_model") or "gemini-3-pro-image-preview"

                        if api_key:
                            log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] Scheduling Gemini Inpainting...")
                            prompt = (
                                "Use the second image as a mask to remove text from the first image. "
                                "Fill the removed text areas with background texture to make it clean. "
                                "Keep non-text areas (figures, tables) unchanged."
                            )
                            
                            async def _run_ai_job(_p_idx=page_idx, _img_p=img_path, _mask_p=str(mask_path), _out_p=str(clean_bg_path)):
                                await _call_image_api_with_retry(
                                    lambda: gemini_multi_image_edit_async(
                                        prompt=prompt,
                                        image_paths=[_img_p, _mask_p],
                                        save_path=_out_p,
                                        api_url=api_url,
                                        api_key=api_key,
                                        model=model_name,
                                        resolution="1K", 
                                        timeout=300
                                    )
                                )
                            
                            ai_task = _run_ai_job()
                            ai_coroutines.append(ai_task)
                        else:
                            log.warning("Skipping AI edit: No API Key provided")
                except Exception as e:
                    log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] Prepare AI task failed: {e}")

            # 保存所有需要在渲染阶段使用的数据
            pages_render_data.append({
                "page_idx": page_idx,
                "scale_x": slide_w_emu / w0,
                "scale_y": slide_h_emu / h0,
                "clean_bg_path": str(clean_bg_path),
                "image_zones": image_zones,
                "final_sam_items": final_sam_items,
                "final_ocr_lines": final_ocr_lines,
                "ai_task": ai_task  # 用于追踪哪个页面发起了 AI 请求
            })

        # ==========================================================
        # Phase 2: 并发执行 AI 任务 & 字号聚类
        # ==========================================================
        
        # 2.1 字号聚类逻辑
        use_global_clustering = getattr(state, "use_global_font_clustering", False)
        global_clusterer = None
        
        if use_global_clustering:
            log.info("[pdf2ppt_with_sam] Performing GLOBAL font size clustering...")
            all_sizes = []
            for p_data in pages_render_data:
                for line in p_data["final_ocr_lines"]:
                    # line: (bbox, text, conf, type, raw_pt)
                    raw_pt = line[4]
                    if raw_pt and raw_pt > 0:
                        all_sizes.append(raw_pt)
            
            global_clusterer = ppt_tool.FontSizeClustering(n_clusters=3)
            global_clusterer.fit(all_sizes)
        
        # 2.2 执行 AI 任务
        if ai_coroutines:
            log.info(f"[pdf2ppt_with_sam] Executing {len(ai_coroutines)} AI background tasks in parallel...")
            start_t = __import__("time").time()
            # 忽略异常，确保后续 PPT 渲染能继续（失败的会降级为白底）
            await asyncio.gather(*ai_coroutines, return_exceptions=True)
            cost = __import__("time").time() - start_t
            log.info(f"[pdf2ppt_with_sam] AI tasks finished. cost={cost:.2f}s")

        # ==========================================================
        # Phase 3: 生成 PPT 页面 (组装)
        # ==========================================================
        for p_data in pages_render_data:
            # 取出数据
            scale_x = p_data["scale_x"]
            scale_y = p_data["scale_y"]
            clean_bg_path = p_data["clean_bg_path"]
            image_zones = p_data["image_zones"]
            final_sam_items = p_data["final_sam_items"]
            final_ocr_lines = p_data["final_ocr_lines"]
            
            # 准备当页的字号聚类器
            if use_global_clustering:
                clusterer = global_clusterer
            else:
                # 单页聚类模式
                page_sizes = [l[4] for l in final_ocr_lines if l[4] > 0]
                clusterer = ppt_tool.FontSizeClustering(n_clusters=3)
                clusterer.fit(page_sizes)

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 3.1 设置背景
            bg_image_path_for_ppt = None
            if os.path.exists(clean_bg_path):
                bg_image_path_for_ppt = clean_bg_path
            
            if bg_image_path_for_ppt:
                try:
                    slide.shapes.add_picture(bg_image_path_for_ppt, 0, 0, prs.slide_width, prs.slide_height)
                except Exception as e:
                    log.error(f"Failed to set slide background image: {e}")
                    # 降级
                    bg = slide.background
                    fill = bg.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(255, 255, 255)
            else:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)

            # 3.2 渲染 MinerU Image Zones
            for zone in image_zones:
                ipath = zone["img_path"]
                if not os.path.exists(ipath): 
                    log.warning(f"MinerU image path not found: {ipath}")
                    continue
                
                bbox = zone["bbox"]
                left = ppt_tool.px_to_emu(bbox[0], scale_x)
                top = ppt_tool.px_to_emu(bbox[1], scale_y)
                width = ppt_tool.px_to_emu(bbox[2] - bbox[0], scale_x)
                height = ppt_tool.px_to_emu(bbox[3] - bbox[1], scale_y)
                
                try:
                    slide.shapes.add_picture(ipath, left, top, width, height)
                except Exception as e:
                    log.error(f"Failed to add mineru image: {e}")

            # 3.3 渲染 SAM Icons
            for item in final_sam_items:
                ipath = item.get("fg_png_path") or item.get("png_path")
                if not ipath or not os.path.exists(ipath): continue
                
                bbox = item.get("bbox_px")
                left = ppt_tool.px_to_emu(bbox[0], scale_x)
                top = ppt_tool.px_to_emu(bbox[1], scale_y)
                width = ppt_tool.px_to_emu(bbox[2] - bbox[0], scale_x)
                height = ppt_tool.px_to_emu(bbox[3] - bbox[1], scale_y)
                
                try:
                    slide.shapes.add_picture(ipath, left, top, width, height)
                except Exception as e:
                    log.error(f"Failed to add SAM icon: {e}")

            # 3.4 渲染 OCR Text
            for line in final_ocr_lines:
                bbox, text, conf, l_type, raw_pt = line
                x1, y1, x2, y2 = bbox
                if (x2 - x1) < 5 or (y2 - y1) < 5: continue

                left = ppt_tool.px_to_emu(x1, scale_x)
                top = ppt_tool.px_to_emu(y1, scale_y)
                width = max(1, ppt_tool.px_to_emu(x2 - x1, scale_x))
                height = max(1, ppt_tool.px_to_emu(y2 - y1, scale_y))

                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = True
                tb.fill.background()
                tb.line.fill.background()

                p = tf.paragraphs[0]
                p.text = text
                
                # 应用字号映射
                final_pt = clusterer.map(raw_pt)
                p.font.size = Pt(final_pt)
                
                # MinerU 的 Title 标签只用于加粗，不再强制改变字号
                if l_type == "title":
                    p.font.bold = True
                
                p.font.color.rgb = RGBColor(0, 0, 0)

        # Save
        # base_dir 已在函数开头定义
        ppt_path = base_dir / "pdf2ppt_with_sam_output.pptx"
        prs.save(str(ppt_path))
        state.ppt_path = str(ppt_path)
        log.info(f"[pdf2ppt_with_sam] PPT generated: {ppt_path}")

        return state

    nodes = {
        "_start_": _init_result_path,
        "pdf_to_images": pdf_to_images_node,
        "parallel_processing": parallel_processing_node,  # 新增：并行处理节点
        "slides_ppt_generation": slides_ppt_generation_node,
        "_end_": lambda state: state,
    }

    edges = [
        ("pdf_to_images", "parallel_processing"),       # pdf_to_images 后进入并行处理
        ("parallel_processing", "slides_ppt_generation"),  # 并行完成后汇合到 PPT 生成
        ("slides_ppt_generation", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges)
    builder.add_edge("_start_", "pdf_to_images")
    return builder
