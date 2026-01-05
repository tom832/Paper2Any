"""
pdf2ppt_with_sam workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
基于 slides PDF:
1. 将 PDF 每页渲染为 PNG
2. 对每页图片用 PaddleOCR 做文字 OCR（沿用 ppt_tool 里的逻辑）
3. 对每页图片用 SAM 做图标 / 图块分割
4. 合并文字 + 图像元素，生成可编辑 PPT

注意：
- 文字完全走 PaddleOCR，不依赖 MinerU；
- 图标 / shape 的分割依然用 SAM；
- 主要用于“PDF版PPT → 可编辑PPT”的场景。
"""

from __future__ import annotations
import os
from pathlib import Path
from typing import List, Dict, Any

import cv2
import fitz  # PyMuPDF
from PIL import Image

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.state import Paper2FigureState 
from dataflow_agent.utils import get_project_root

# SAM & PPT helpers
from dataflow_agent.toolkits.imtool.sam_tool import segment_layout_boxes, free_sam_model
from dataflow_agent.toolkits.imtool.bg_tool import local_tool_for_bg_remove, free_bg_rm_model
from dataflow_agent.toolkits.imtool import ppt_tool

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

log = get_logger(__name__)


def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    为本次 pdf2ppt_with_sam workflow 创建统一的输出目录：
    - 如果 state.result_path 已存在，直接使用；
    - 否则使用项目根目录下 outputs/pdf2ppt_with_sam/<timestamp>。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(__import__("time").time())
    base_dir = (root / "outputs" / "pdf2ppt_with_sam" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path


def _run_sam_on_pages(image_paths: List[str], base_dir: str) -> List[Dict[str, Any]]:
    """
    对每一页图片运行 SAM，输出 layout_items：
    - 每页一个 dict: {"page_idx": i, "layout_items": [...layout_item...]}
    - layout_item 内部含:
        - bbox (norm)
        - bbox_px (像素坐标)
        - png_path: SAM 裁剪出的原始小图（PNG）
        - fg_png_path: 可选，后续背景抠图后的 PNG
    """
    results: List[Dict[str, Any]] = []
    sam_ckpt = f"{get_project_root()}/sam_b.pt"

    for page_idx, img_path in enumerate(image_paths):
        img_path_obj = Path(img_path)
        if not img_path_obj.exists():
            log.warning(f"[pdf2ppt_with_sam] image not found for SAM: {img_path}")
            results.append({"page_idx": page_idx, "layout_items": []})
            continue

        out_dir = Path(base_dir) / "layout_items" / f"page_{page_idx+1:03d}"
        out_dir.mkdir(parents=True, exist_ok=True)

        # 1. SAM 分割
        layout_items = segment_layout_boxes(
            image_path=str(img_path_obj),
            output_dir=str(out_dir),
            checkpoint=sam_ckpt,
            min_area=200,
            min_score=0.0,
            iou_threshold=0.2,
            top_k=15,
            nms_by="mask",
        )
        log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] SAM found {len(layout_items)} items")

        # 2. 映射 bbox 到像素坐标（基于整页尺寸）
        try:
            pil_img = Image.open(str(img_path_obj))
            w, h = pil_img.size
        except Exception as e:
            log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] open image failed: {e}")
            w, h = 1024, 768

        for it in layout_items:
            bbox = it.get("bbox")
            if bbox and len(bbox) == 4:
                x1n, y1n, x2n, y2n = bbox
                x1 = int(round(x1n * w))
                y1 = int(round(y1n * h))
                x2 = int(round(x2n * w))
                y2 = int(round(y2n * h))
                if x2 > x1 and y2 > y1:
                    it["bbox_px"] = [x1, y1, x2, y2]

        results.append({"page_idx": page_idx, "layout_items": layout_items})

    # 显式释放 SAM 模型
    try:
        free_sam_model(checkpoint=sam_ckpt)
    except Exception as e:
        log.error(f"[pdf2ppt_with_sam] free_sam_model failed: {e}")

    return results


# 不再使用 EMF，保留占位函数（如需调试可删除）
def _add_layout_emf_to_slide(*args, **kwargs):  # pragma: no cover
    return False


@register("pdf2ppt_with_sam")
def create_pdf2ppt_with_sam_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf pdf2ppt_with_sam

    要求:
    - state.pdf_file: 输入 PDF 路径（slides）
    - 可选: state.result_path: 输出根目录，如未提供自动生成
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState, entry_point="_start_")

    # ==============================
    # NODES
    # ==============================

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        _ensure_result_path(state)
        return state

    async def pdf_to_images_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        将 PDF 每一页渲染为 PNG。
        """
        pdf_path = getattr(state, "pdf_file", None)
        if not pdf_path:
            log.error("[pdf2ppt_with_sam] state.pdf_file is empty")
            return state

        base_dir = Path(_ensure_result_path(state))
        img_dir = base_dir / "slides_png"
        image_paths = ppt_tool.pdf_to_images(pdf_path, str(img_dir))
        state.slide_images = image_paths
        return state

    async def slides_ocr_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        对每一页图片用 PaddleOCR 做 OCR。
        复用 ppt_tool 中的逻辑，输出结构化的 OCR 行信息。
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.error("[pdf2ppt_with_sam] no slide_images for OCR")
            return state

        ocr_pages: List[Dict[str, Any]] = []
        for page_idx, img_path in enumerate(image_paths):
            try:
                result = ppt_tool.paddle_ocr_page_with_layout(img_path)
            except Exception as e:
                log.error(f"[pdf2ppt_with_sam][OCR] page#{page_idx+1} failed: {e}")
                result = {
                    "image_size": None,
                    "lines": [],
                    "body_h_px": None,
                    "bg_color": None,
                    "path": img_path,
                    "page_idx": page_idx,
                }
            result["page_idx"] = page_idx
            result["path"] = img_path
            ocr_pages.append(result)

        state.ocr_pages = ocr_pages
        return state

    async def slides_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        对每一页图片运行 SAM 用于图标 / 图块分割。
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.error("[pdf2ppt_with_sam] no slide_images for SAM")
            return state

        base_dir = _ensure_result_path(state)
        sam_pages = _run_sam_on_pages(image_paths, base_dir)
        state.sam_pages = sam_pages
        return state

    async def slides_layout_bg_remove_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        对每一页 SAM layout PNG 做背景抠图：
        - 输入: state.sam_pages[*].layout_items[].png_path
        - 输出: 为每个 layout_item 写入 fg_png_path（抠完背景的 PNG）
        """
        sam_pages: List[Dict[str, Any]] = getattr(state, "sam_pages", []) or []
        if not sam_pages:
            log.error("[pdf2ppt_with_sam] no sam_pages for bg remove")
            return state

        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "sam_icons"
        icons_dir.mkdir(parents=True, exist_ok=True)

        model_path = getattr(getattr(state, "request", None), "bg_rm_model", None)

        processed = 0
        
        for p in sam_pages:
            page_idx = p.get("page_idx", 0)
            for it in p.get("layout_items", []):
                png_path = it.get("png_path")
                if not png_path or not os.path.exists(png_path):
                    continue
                
                # 背景抠图 - 添加页码前缀避免文件名冲突
                try:
                    # 从原始路径提取文件名
                    original_stem = Path(png_path).stem
                    # 创建带页码的输出文件名
                    output_filename = f"page_{page_idx+1:03d}_{original_stem}_bg_removed.png"
                    output_path = icons_dir / output_filename
                    
                    req = {
                        "image_path": png_path,
                        "output_dir": str(icons_dir),
                    }
                    if model_path:
                        req["model_path"] = model_path
                    
                    fg_path = local_tool_for_bg_remove(req)
                    
                    # 重命名文件以包含页码
                    if fg_path and os.path.exists(fg_path):
                        # 将生成的文件重命名为带页码的文件名
                        fg_path_obj = Path(fg_path)
                        if fg_path_obj.name != output_filename:
                            new_fg_path = fg_path_obj.parent / output_filename
                            fg_path_obj.rename(new_fg_path)
                            fg_path = str(new_fg_path)
                        
                        it["fg_png_path"] = fg_path
                    else:
                        it["fg_png_path"] = png_path
                    
                    processed += 1
                except Exception as e:
                    log.error(f"[pdf2ppt_with_sam][bg_rm] failed for {png_path}: {e}")
                    it["fg_png_path"] = png_path

        # 抠图完成后可尝试释放模型（忽略失败）
        try:
            if model_path:
                free_bg_rm_model(model_path=model_path)
        except Exception as e:
            log.error(f"[pdf2ppt_with_sam] free_bg_rm_model failed: {e}")

        log.info(f"[pdf2ppt_with_sam] bg remove processed: {processed} items")
        return state

    async def slides_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        结合 OCR + SAM 结果生成可编辑 PPT：
        - 背景: 使用图像编辑API进行inpainting后的干净底图
        - 文本: PaddleOCR 行结果，使用 ppt_tool 的字号/颜色逻辑
        - 布局形状: 使用 SAM 生成的 PNG 图标
        """
        from dataflow_agent.toolkits.imtool.req_img import generate_or_edit_and_save_image_async
        
        ocr_pages: List[Dict[str, Any]] = getattr(state, "ocr_pages", []) or []
        sam_pages: List[Dict[str, Any]] = getattr(state, "sam_pages", []) or []

        if not ocr_pages:
            log.error("[pdf2ppt_with_sam] no ocr_pages, abort PPT generation")
            return state

        # 以 PPT 工具里的默认比例创建 Presentation
        prs = Presentation()
        prs.slide_width = Inches(ppt_tool.SLIDE_W_IN)
        prs.slide_height = Inches(ppt_tool.SLIDE_H_IN)

        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height

        # 简单做一个 page_idx 到 sam_items 的映射
        sam_dict: Dict[int, List[Dict[str, Any]]] = {}
        for p in sam_pages:
            page_idx = p.get("page_idx", 0)
            sam_dict[page_idx] = p.get("layout_items", [])

        def _filter_sam_items_for_ppt(
            items: List[Dict[str, Any]],
            img_w: int,
            img_h: int,
            ocr_lines: List[Any],
        ) -> List[Dict[str, Any]]:
            """
            二次过滤 SAM 图块，适配 PPT 视觉需求：
            - 根据整页面积过滤过小/过大的块（噪声 / 整块背景）
            - 根据宽高比过滤极端细长的块（大概率是线条而非 icon）
            - 可选：与 OCR 文本框高度重叠的块可丢弃，避免把整段文字当图标
            """
            if not items:
                return []

            total_area = img_w * img_h

            # 预先取出文本 bbox 列表，供简单 IOU 过滤使用
            text_bboxes = [b for (b, _t, _c) in ocr_lines] if ocr_lines else []

            def _bbox_iou(a, b) -> float:
                ax1, ay1, ax2, ay2 = a
                bx1, by1, bx2, by2 = b
                ix1, iy1 = max(ax1, bx1), max(ay1, by1)
                ix2, iy2 = min(ax2, bx2), min(ay2, by2)
                iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
                inter = iw * ih
                if inter <= 0:
                    return 0.0
                area_a = (ax2 - ax1) * (ay2 - ay1)
                area_b = (bx2 - bx1) * (by2 - by1)
                return inter / (area_a + area_b - inter + 1e-6)

            filtered: List[Dict[str, Any]] = []
            for it in items:
                bbox = it.get("bbox_px") or it.get("bbox")
                if not bbox or len(bbox) != 4:
                    continue
                x1, y1, x2, y2 = bbox
                w = max(1, x2 - x1)
                h = max(1, y2 - y1)
                area = w * h

                # 面积比例过滤：去掉极小噪声块和特别大的背景块
                if area < 0.0005 * total_area:  # <0.05%
                    continue
                if area > 0.3 * total_area:  # >30% 整页，通常不是 icon
                    continue

                # 形状过滤：太细长的大概率是线条，不当 icon
                aspect = max(w / h, h / w)
                if aspect > 8.0:
                    continue

                # 可选：如果与某个文本框高度重叠（IOU 高），认为这是文字区域，丢掉
                drop_for_text_iou = False
                for tb in text_bboxes:
                    if _bbox_iou(bbox, tb) > 0.8:
                        drop_for_text_iou = True
                        break
                if drop_for_text_iou:
                    continue

                filtered.append(it)
            return filtered

        for pinfo in ocr_pages:
            page_idx = pinfo.get("page_idx", 0)
            img_path = pinfo.get("path")
            lines = pinfo.get("lines", [])
            body_h_px = pinfo.get("body_h_px")
            bg_color = pinfo.get("bg_color")
            if not img_path or not os.path.exists(img_path):
                log.warning(f"[pdf2ppt_with_sam] missing img for page#{page_idx+1}: {img_path}")
                continue

            bgr = ppt_tool.read_bgr(img_path)
            h0, w0 = bgr.shape[:2]

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # 背景填白
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            # 1) 使用图像编辑API进行inpainting生成干净底图（带重试机制）
            clean_bg = bgr  # 默认使用原图
            if ppt_tool.ADD_BACKGROUND_IMAGE and lines:
                import asyncio
                
                async def _call_inpaint_api_with_retry(retries: int = 3, delay: float = 1.0) -> bool:
                    """
                    对 inpainting API 进行最多 retries 次重试
                    - 成功：返回 True
                    - 多次失败：返回 False
                    """
                    last_err: Optional[Exception] = None
                    for attempt in range(1, retries + 1):
                        try:
                            await generate_or_edit_and_save_image_async(
                                prompt=inpaint_prompt,
                                save_path=str(clean_bg_path),
                                aspect_ratio="16:9",
                                api_url=state.request.chat_api_url,
                                api_key=state.request.api_key or state.request.chat_api_key or os.getenv("DF_API_KEY") ,
                                model=state.request.gen_fig_model,
                                image_path=str(temp_img_path),
                                use_edit=True,
                            )
                            return True
                        except Exception as e:
                            last_err = e
                            log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] inpainting attempt {attempt}/{retries} failed: {e}")
                            if attempt < retries:
                                try:
                                    await asyncio.sleep(delay)
                                except Exception:
                                    pass
                    log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] inpainting failed after {retries} attempts: {last_err}")
                    return False
                
                try:
                    # 生成文字掩码
                    text_mask = ppt_tool.build_adaptive_mask(bgr, lines)
                    
                    # 保存掩码图（用于调试）
                    base_dir = Path(_ensure_result_path(state))
                    mask_dir = base_dir / "text_masks"
                    mask_dir.mkdir(parents=True, exist_ok=True)
                    mask_path = mask_dir / f"text_mask_{page_idx+1:03d}.png"
                    cv2.imwrite(str(mask_path), text_mask)
                    
                    # 保存原图临时文件
                    temp_img_path = mask_dir / f"temp_original_{page_idx+1:03d}.png"
                    cv2.imwrite(str(temp_img_path), bgr)
                    
                    # 构造inpainting提示词
                    inpaint_prompt = "请智能修复图像中文字被移除后的区域，保持背景的连续性、一致性和自然过渡，使修复后的图像看起来完整无缺，并且原图涉及的图标你需要尽量保留；"
                    
                    # 调用图像编辑API进行inpainting（带重试）
                    clean_bg_path = base_dir / "clean_backgrounds" / f"clean_bg_{page_idx+1:03d}.png"
                    clean_bg_path.parent.mkdir(parents=True, exist_ok=True)
                    
                    log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] 开始调用图像编辑API进行inpainting（最多重试3次）...")
                    
                    api_success = await _call_inpaint_api_with_retry(retries=3, delay=1.0)
                    
                    # 读取修复后的图像作为底图
                    if api_success and os.path.exists(str(clean_bg_path)):
                        clean_bg = ppt_tool.read_bgr(str(clean_bg_path))
                        log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] API inpainting成功，使用修复后的底图")
                    else:
                        log.warning(f"[pdf2ppt_with_sam][page#{page_idx+1}] API inpainting失败，使用本地inpaint作为fallback")
                        # Fallback: 使用本地OpenCV inpaint
                        try:
                            clean_bg = ppt_tool.make_clean_background(bgr, lines)
                            log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] 使用本地inpaint成功")
                        except Exception as e2:
                            log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] 本地inpaint也失败: {e2}，使用原图")
                            clean_bg = bgr
                    
                    # 清理临时文件
                    if os.path.exists(str(temp_img_path)):
                        os.remove(str(temp_img_path))
                        
                except Exception as e:
                    log.error(f"[pdf2ppt_with_sam][page#{page_idx+1}] inpainting整体流程失败: {e}，使用原图")
                    clean_bg = bgr
                
                tmp = f"__pdf2ppt_bg_{page_idx+1}.png"
                ppt_tool.add_background(slide, clean_bg, slide_w_emu, slide_h_emu, tmp)

            # 2) 像素→emu 缩放比例（SAM 图层 + 文本共用）
            scale_x = slide_w_emu / w0
            scale_y = slide_h_emu / h0

            # 3) 叠加 SAM 的 PNG 图形（抠图后的图标）
            sam_items = sam_dict.get(page_idx, [])
            sam_items = _filter_sam_items_for_ppt(sam_items, w0, h0, lines)
            layout_drawn = 0
            
            for item in sam_items:
                fg_path = item.get("fg_png_path") or item.get("png_path")
                bbox = item.get("bbox_px") or item.get("bbox")
                
                if not fg_path or not bbox or len(bbox) != 4:
                    continue

                x1, y1, x2, y2 = bbox
                if x2 <= x1 or y2 <= y1:
                    continue

                left = ppt_tool.px_to_emu(x1, scale_x)
                top = ppt_tool.px_to_emu(y1, scale_y)
                width = max(1, ppt_tool.px_to_emu((x2 - x1), scale_x))
                height = max(1, ppt_tool.px_to_emu((y2 - y1), scale_y))

                # 使用抠图后的 PNG
                if os.path.exists(fg_path):
                    try:
                        slide.shapes.add_picture(fg_path, left, top, width, height)
                        layout_drawn += 1
                    except Exception as e:
                        log.error(f"[pdf2ppt_with_sam] add PNG layout failed: {e}")

            log.info(f"[pdf2ppt_with_sam][page#{page_idx+1}] SAM layout PNG drawn: {layout_drawn}")

            # 4) 绘制文字文本框（完全沿用 ppt_tool 的计算方式）

            for bbox, text, conf in lines:
                x1, y1, x2, y2 = bbox
                if (x2 - x1) < 6 or (y2 - y1) < 6:
                    continue

                left = ppt_tool.px_to_emu(x1, scale_x)
                top = ppt_tool.px_to_emu(y1, scale_y)
                width = max(1, ppt_tool.px_to_emu((x2 - x1), scale_x))
                height = max(1, ppt_tool.px_to_emu((y2 - y1), scale_y))

                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = True

                tb.fill.background()
                tb.line.fill.background()

                p = tf.paragraphs[0]
                p.text = text
                p.font.size = ppt_tool.estimate_font_pt(
                    bbox, img_h_px=h0, body_h_px=body_h_px
                )

                if ppt_tool.EXTRACT_TEXT_COLOR:
                    text_color = ppt_tool.extract_text_color(bgr, bbox, bg_color)
                    p.font.color.rgb = RGBColor(*text_color)
                else:
                    p.font.color.rgb = RGBColor(0, 0, 0)

        base_dir = Path(_ensure_result_path(state))
        ppt_path = base_dir / "pdf2ppt_with_sam_output.pptx"
        prs.save(str(ppt_path))
        state.ppt_path = str(ppt_path)
        log.info(f"[pdf2ppt_with_sam] PPT generated: {ppt_path}")

        return state

    nodes = {
        "_start_": _init_result_path,
        "pdf_to_images": pdf_to_images_node,
        "slides_ocr": slides_ocr_node,
        "slides_sam": slides_sam_node,
        "slides_layout_bg_remove": slides_layout_bg_remove_node,
        "slides_ppt_generation": slides_ppt_generation_node,
        "_end_": lambda state: state,
    }

    edges = [
        ("pdf_to_images", "slides_ocr"),
        ("slides_ocr", "slides_sam"),
        ("slides_sam", "slides_layout_bg_remove"),
        ("slides_layout_bg_remove", "slides_ppt_generation"),
        ("slides_ppt_generation", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges)
    # 入口从 _start_ 固定到 pdf_to_images
    builder.add_edge("_start_", "pdf_to_images")
    return builder
