"""
icongen workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
生成时间: 2025-10-27 11:11:56

1. 在 **TOOLS** 区域定义需要暴露给 Prompt 的前置工具
2. 在 **NODES**  区域实现异步节点函数 (await-able)
3. 在 **EDGES**  区域声明有向边
4. 最后返回 builder.compile() 或 GenericGraphBuilder
"""

from __future__ import annotations
import asyncio
import json
import os
from dataflow_agent.state import MainState, Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder


from dataflow_agent.workflow.registry import register
# from dataflow_agent.agentroles import get_agent_cls, create_agent

from dataflow_agent.toolkits.tool_manager import get_tool_manager
from langchain.tools import tool
from langgraph.graph import StateGraph
from langgraph.prebuilt import ToolNode, tools_condition

from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.toolkits.imtool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.toolkits.imtool.bg_tool import local_tool_for_bg_remove, local_tool_for_raster_to_svg, free_bg_rm_model
from dataflow_agent.toolkits.imtool.sam_tool import segment_layout_boxes, segment_layout_boxes_server, free_sam_model
from dataflow_agent.toolkits.imtool.mineru_tool import (
    svg_to_emf,
    recursive_mineru_layout,
)
from dataflow_agent.agentroles import create_graph_agent,create_react_agent

import re, pdfplumber, PyPDF2, time, shutil, fitz
import numpy as np
from PIL import Image

from dataflow_agent.utils import (
    build_output_directory,
    add_image_element,
    add_text_element,
    setup_presentation_size,
    get_project_root,
    pixels_to_inches,
)

from pathlib import Path
import time, random
from pptx import Presentation
from pptx.dml.color import RGBColor 
from pptx.util import Inches


log = get_logger(__name__)

TEMPLATE_EDIT_PROMPT = (
"Transform the original image into a pure layout made ONLY of solid colored blocks:\n"
"1. Keep only the outermost rectangles and arrows (if they exist).\n"
"2. Delete everything inside them: all titles, subtitles, texts, icons, illustrations, and any inner shapes.\n"
"3. Turn each remaining outer shape into a solid color block; remove borders if possible.\n"
"4. Keep the layout exactly the same: same positions, sizes, alignment, and spacing.\n"
"5. Do NOT add any text, labels, or symbols anywhere.\n"
"Finally, output a description of this empty color-block template (no text content at all)."
)

def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    统一本次 paper2figure_with_sam workflow 的根输出目录：
    - 如果 state.result_path 已存在（通常由调用方传入），直接使用；
    - 否则：使用 get_project_root() / "outputs" / "paper2figure" / <timestamp>，
      并写回 state.result_path，后续节点共享同一目录。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(time.time())
    base_dir = (root / "outputs" / "paper2figure" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path

def _ts_name(stem: str, ext: str = ".png") -> str:
    timestamp = int(time.time())  # 获取当前时间戳（秒）
    return f"./{stem}{timestamp}{ext}"

@register("paper2fig_with_sam")
def create_p2fig_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf paper2fig
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState,
                                  entry_point="_start_")  # 自行修改入口

    # ----------------------------------------------------------------------
    # TOOLS (pre_tool definitions)
    # ----------------------------------------------------------------------
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_abstract_intro(state: Paper2FigureState):
        """
        Robustly extract Abstract + Introduction from PDF.
        """

        # 1. Read metadata title
        try:
            with open(state.paper_file, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                paper_title = reader.metadata.get('/Title', 'Unknown Title')
        except Exception:
            paper_title = "Unknown Title"

        # Open the PDF file using the path from state
        file_path = state.paper_file
        pdf_document = fitz.open(file_path)

        # Extract text from the first 10 pages
        text = ""
        for page_num in range(min(10, len(pdf_document))):
            page = pdf_document.load_page(page_num)
            text += page.get_text("text")

        content = text.strip()

        final_text = (
            f"The title of the paper is {paper_title}\n\n"
            f"Here's first ten page content: {content}"
        )

        log.info(f"{final_text}")
        return final_text
    
    @builder.pre_tool("paper_idea", "figure_desc_generator")
    def _get_paper_idea(state: Paper2FigureState):
        """
        Return paper ideas summary.
        """
        return state.paper_idea

    # ==============================================================
    # NODES
    # ==============================================================
    async def paper_idea_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        提取论文的关键贡献点
        """
        paper_idea_extractor = create_graph_agent("paper_idea_extractor")
        state = await paper_idea_extractor.execute(state, use_agent=True)
        return state
    
    async def figure_desc_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        图标提示词生成器节点
        """
        figure_desc_generator = create_react_agent("figure_desc_generator",
                                                    max_retries=5,
                                                    model_name="gpt-5.1")
        state = await figure_desc_generator.execute(state, use_agent=True)
        return state

    async def figure_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        图像生成或编辑节点：
        1) 先生成带内容的图 (fig_draft_path)
        2) 再基于该图进行一次固定提示词的二次编辑，生成空框模板图 (fig_layout_path)
        """
        prompt = state.agent_results.get("figure_desc_generator").get("results").get("fig_desc", {})
        safe_prompt = json.dumps(prompt, ensure_ascii=False)  # 确保中文字符正常显示

        edit_prompt = state.request.get("edit_prompt")
        image_path = state.request.get("prev_image")

        # 如果是二次编辑，prompt可以为空
        final_prompt = edit_prompt if image_path else safe_prompt

        log.info(f'final_prompt{final_prompt} - edit_prompt：{edit_prompt} - image_path：{image_path} - prompt：{safe_prompt}')

        # 统一输出根目录（outputs/paper2figure/<ts>）
        result_root = Path(_ensure_result_path(state))
        result_root.mkdir(parents=True, exist_ok=True)

        # 1) 生成带内容的图，直接存到 result_root
        fig_name = f"fig_{int(time.time())}.png"
        save_path = str(result_root / fig_name)

        await generate_or_edit_and_save_image_async(
            prompt=final_prompt,
            save_path=save_path,
            aspect_ratio=state.aspect_ratio,
            api_url=state.request.chat_api_url,
            api_key=state.request.chat_api_key or os.getenv("DF_API_KEY") ,
            model=state.request.gen_fig_model,
            image_path=image_path,
            use_edit=True if image_path else False
        )
        state.agent_results["gen_img"] = {"path": save_path}
        state.fig_draft_path = save_path

        # 2) 基于第一次生成的图，做一次“空模板”二次编辑，也放在 result_root
        # TEMPLATE_EDIT_PROMPT = (
        #     "Keep only the outermost rectangles and arrows(if any in the original box).\n"
        #     "Remove all inner content including title, subtitles, icons, explainary texts and all that.\n"
        #     "Keep the layout exactly the same.\n"
        #     "Output a description of an empty template composed of these boxes."
        # )

        layout_name = f"layout_{int(time.time())}.png"
        layout_save_path = str(result_root / layout_name)
        await generate_or_edit_and_save_image_async(
            prompt=TEMPLATE_EDIT_PROMPT,
            save_path=layout_save_path,
            aspect_ratio=state.aspect_ratio,
            api_url=state.request.chat_api_url,
            api_key=state.request.chat_api_key or os.getenv("DF_API_KEY") ,
            model=state.request.gen_fig_model,
            image_path=save_path,
            use_edit=True,
        )
        state.fig_layout_path = layout_save_path
        state.agent_results["gen_img_template"] = {"path": layout_save_path}

        return state

    async def figure_layout_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        针对二次编辑后的空框模板图 (fig_layout_path) 进行:
        SAM 自动分割 -> 过滤 -> 裁剪子图 -> PNG->SVG->EMF，
        结果写入 state.layout_items，仅作为 PPT 背景框架层。

        注意：
        - segment_layout_boxes 返回的 bbox 是基于 layout 图像尺寸的归一化坐标 [0,1]；
        - 这里显式转换一份像素坐标 bbox_px，后面插入 PPT 时统一按像素 → 英寸 → Emu 的规则处理，
          和 add_image_element / add_text_element 的坐标系保持一致，避免 EMF 位置/尺寸错乱导致“看不到”的问题。
        """
        if not state.fig_layout_path and state.request.input_type == "FIGURE":
            result_root = Path(_ensure_result_path(state))
            result_root.mkdir(parents=True, exist_ok=True)
            log.critical(f"[figure_layout_sam] fig_layout_path 为空， 需要更新Layout图")
            layout_name = f"layout_{int(time.time())}.png"
            layout_save_path = str(result_root / layout_name)
            await generate_or_edit_and_save_image_async(
                prompt="1.Remove all text content; keep only the outermost rectangular frames and arrows (if any).\n"
                       "2.Keep the layout unchanged.\n"
                       "3.Change the background color to white.",
                save_path=layout_save_path,
                aspect_ratio=state.aspect_ratio,
                api_url=state.request.chat_api_url,
                api_key=state.request.chat_api_key or os.getenv("DF_API_KEY") ,
                model=state.request.gen_fig_model,
                image_path=f"{get_project_root()}/{state.fig_draft_path}",
                use_edit=True,
            )
            state.fig_layout_path = layout_save_path
            state.agent_results["gen_img_template"] = {"path": layout_save_path}

        img_path = Path(state.fig_layout_path)
        if not img_path.exists():
            log.error(f"[figure_layout_sam] fig_layout_path 不存在: {img_path}")
            return state

        base_dir = Path(_ensure_result_path(state))
        out_dir = base_dir / "layout_items"
        out_dir.mkdir(parents=True, exist_ok=True)

        sam_ckpt = f'{get_project_root()}/sam_b.pt'
        # SAM LB Port 8020
        sam_server_urls = ["http://localhost:8020"]

        # 1. SAM 分割 + 过滤 + 裁剪子图 (优先使用远程服务)
        try:
            layout_items = segment_layout_boxes_server(
                image_path      = str(img_path),
                output_dir      = str(out_dir),
                server_urls     = sam_server_urls,
                checkpoint      = sam_ckpt,
                min_area        = 200,
                min_score       = 0.0,
                iou_threshold   = 0.2,
                top_k           = 15,
                nms_by          = "mask",
            )
        except Exception as e:
            log.error(f"[figure_layout_sam] Remote SAM failed: {e}. Fallback to local.")
            # Fallback to local if server fails
            layout_items = segment_layout_boxes(
                image_path      = str(img_path),
                output_dir      = str(out_dir),
                checkpoint      = sam_ckpt,
                # 这里的参数可以根据 mask_detail_level 调整
                min_area        = 200,
                min_score       = 0.0,
                iou_threshold   = 0.2,
                top_k           = 15,
                nms_by          = "mask",
            )
            # 只有本地运行时才需要手动释放模型
            free_sam_model(checkpoint= sam_ckpt)

        log.info(f"[figure_layout_sam] SAM 分割结果: {len(layout_items)} 个布局元素")

        # layout 图实际像素尺寸，用于把归一化 bbox 转为像素 bbox
        try:
            layout_img = Image.open(str(img_path))
            layout_w, layout_h = layout_img.size
        except Exception as e:
            log.error(f"[figure_layout_sam] 打开 layout 图失败: {e}")
            layout_w, layout_h = 1024, 1024  # 兜底，和默认 slide 尺寸一致

        # 2. 每个 layout PNG 转 SVG -> EMF，并补充像素坐标 bbox_px
        for idx, it in enumerate(layout_items):
            png_path = it.get("png_path")
            if not png_path:
                continue

            # 将归一化 bbox 映射到像素坐标，和 fig_mask 的像素 bbox 保持一致
            bbox = it.get("bbox")
            if bbox and len(bbox) == 4:
                x1n, y1n, x2n, y2n = bbox
                x1 = int(round(x1n * layout_w))
                y1 = int(round(y1n * layout_h))
                x2 = int(round(x2n * layout_w))
                y2 = int(round(y2n * layout_h))
                if x2 > x1 and y2 > y1:
                    it["bbox_px"] = [x1, y1, x2, y2]
                else:
                    log.warning(f"[figure_layout_sam] 无效 bbox: {bbox} -> 像素 [{x1},{y1},{x2},{y2}]")

            svg_path = out_dir / f"layout_{idx}.svg"
            svg_abs = local_tool_for_raster_to_svg(
                {
                    "image_path": png_path,
                    "output_svg": str(svg_path),
                    "colormode": "color",
                    "hierarchical": "stacked",
                    "mode": "spline",
                }
            )
            it["svg_path"] = svg_abs

            emf_path = out_dir / f"layout_{idx}.emf"
            try:
                emf_abs = svg_to_emf(svg_abs, str(emf_path))
                it["emf_path"] = emf_abs
            except Exception as e:
                log.error(f"[figure_layout_sam] svg_to_emf failed for {svg_abs}: {e}")
                it["emf_path"] = None

        state.layout_items = layout_items
        log.info(f"[figure_layout_sam] 共生成 {len(layout_items)} 个布局元素")
        # log.info(f'state.layout_items : {state.layout_items}')
        return state

    async def figure_mask_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成Figure进行元素切割，并提取 bbox + image_path 信息，递归处理子图。
        使用 MinerU HTTP 对原始带内容的图 (fig_draft_path) 做解析，得到内容层元素。
        规则：
        - 标题块(type == 'title') 保留为 text；
        - 其它所有块一律从顶层图裁剪出子图，当作 image，用于 icon / 局部视觉元素。
        """

        img_path = Path(state.fig_draft_path)
        if not img_path.exists():
            log.error(f"[figure_mask] fig_draft_path 不存在: {img_path}")
            return state

        # MinerU 所有中间结果统一放在本次 outputs 下
        base_dir = Path(_ensure_result_path(state))
        out_dir = base_dir / "mineru_recursive"
        out_dir.mkdir(parents=True, exist_ok=True)
        log.info(f"[figure_mask] MinerU 输出目录: {out_dir}")

        # MinerU 端口：优先从 state.request.mineru_port 读取，默认 8010
        port = getattr(state.request, "mineru_port", 8010)
        max_depth = getattr(state, "mask_detail_level", 3)

        log.critical(f"mask detail level : {max_depth} ")
        log.critical(f'[img_path]: {img_path}')
        log.critical(f'[mineru_port]: {port}')

        # 1. 调用新的 HTTP MinerU 递归处理，获取元素列表（归一化坐标）
        mineru_items = await recursive_mineru_layout(
            image_path=str(img_path),
            port=port,
            max_depth=max_depth,
            output_dir=out_dir,
        )
        log.info(f"mineru_items : {mineru_items}")

        # 顶层图像尺寸，用于 norm->pixel 映射与裁剪
        top_img = Image.open(state.fig_draft_path)
        top_w, top_h = top_img.size

        # 图标原图输出目录
        icons_raw_dir = base_dir / "icons_raw"
        icons_raw_dir.mkdir(parents=True, exist_ok=True)

        fig_mask = []
        icon_count = 0
        text_count = 0

        details = 1
        if state.request.figure_complex == "easy":
            details = 1
        elif state.request.figure_complex == "hard":
            details = 10
        else:
            details = 5

        # 如果 MinerU 只返回了 小于等于 6 个整体元素：按 SAM 布局切子图，再对每个子图单独跑 MinerU，
        # 以便获取该布局块内部的文字和更细粒度元素。底图始终使用 fig_draft_path。
        if len(mineru_items) <= details:
            from dataflow_agent.toolkits.imtool.mineru_tool import run_aio_two_step_extract

            layout_items = getattr(state, "layout_items", None) or []
            log.info(f"[figure_mask] mineru_items size = {len(mineru_items)}, 使用 SAM 布局({len(layout_items)} 个)进行二次 MinerU 拆分")

            # 子图保存目录
            sub_root_dir = base_dir / "mineru_sub_images"
            sub_root_dir.mkdir(parents=True, exist_ok=True)

            for layout_idx, layout_it in enumerate(layout_items):
                # 优先使用在 figure_layout_sam_node 中写入的像素 bbox_px；否则退回 bbox 视为像素坐标
                bbox_px = layout_it.get("bbox_px") or layout_it.get("bbox")
                if not bbox_px or len(bbox_px) != 4:
                    continue
                lx1, ly1, lx2, ly2 = bbox_px
                # 粗筛 bbox
                if lx2 <= lx1 or ly2 <= ly1:
                    continue

                # 边界裁剪到顶层图尺寸
                lx1 = max(0, min(top_w, int(round(lx1))))
                ly1 = max(0, min(top_h, int(round(ly1))))
                lx2 = max(0, min(top_w, int(round(lx2))))
                ly2 = max(0, min(top_h, int(round(ly2))))
                if lx2 <= lx1 or ly2 <= ly1:
                    continue

                # 1) 从原始 fig_draft_path 裁出当前布局块子图
                try:
                    sub_img = top_img.crop((lx1, ly1, lx2, ly2))
                except Exception as e:
                    log.error(f"[figure_mask] 裁剪 SAM 子图失败 layout_idx={layout_idx}, bbox=({lx1},{ly1},{lx2},{ly2}): {e}")
                    continue

                sub_dir = sub_root_dir / f"layout_{layout_idx}"
                sub_dir.mkdir(parents=True, exist_ok=True)
                sub_path = sub_dir / f"sam_sub_{layout_idx}.png"
                try:
                    sub_img.save(sub_path)
                except Exception as e:
                    log.error(f"[figure_mask] 保存 SAM 子图失败 layout_idx={layout_idx}, path={sub_path}: {e}")
                    continue

                # 2) 对子图再次调用 MinerU（只做一层 two_step_extract，不再递归）
                try:
                    sub_blocks = await run_aio_two_step_extract(str(sub_path), port=port)
                except Exception as e:
                    log.error(f"[figure_mask] 子图 MinerU 解析失败 layout_idx={layout_idx}, path={sub_path}: {e}")
                    continue

                sub_w, sub_h = sub_img.size

                # 3) 遍历子图内的 MinerU block，映射到整图像素坐标系
                for blk_idx, blk in enumerate(sub_blocks):
                    blk_type_raw = blk.get("type") or ""
                    blk_type = blk_type_raw.lower()
                    bbox_norm = blk.get("bbox")
                    text = blk.get("text") or blk.get("content") or ""
                    if not bbox_norm or len(bbox_norm) != 4:
                        continue

                    sx1n, sy1n, sx2n, sy2n = bbox_norm
                    # 规整到 [0,1]，避免越界
                    sx1n = max(0.0, min(1.0, float(sx1n)))
                    sy1n = max(0.0, min(1.0, float(sy1n)))
                    sx2n = max(0.0, min(1.0, float(sx2n)))
                    sy2n = max(0.0, min(1.0, float(sy2n)))
                    if sx2n <= sx1n or sy2n <= sy1n:
                        continue

                    # 子图归一化 -> 子图像素
                    sx1 = int(round(sx1n * sub_w))
                    sy1 = int(round(sy1n * sub_h))
                    sx2 = int(round(sx2n * sub_w))
                    sy2 = int(round(sy2n * sub_h))
                    if sx2 <= sx1 or sy2 <= sy1:
                        continue

                    # 子图像素 -> 整图像素（加上布局块的偏移）
                    gx1 = lx1 + sx1
                    gy1 = ly1 + sy1
                    gx2 = lx1 + sx2
                    gy2 = ly1 + sy2

                    # 再次 clamp 到整图范围
                    gx1 = max(0, min(top_w, gx1))
                    gy1 = max(0, min(top_h, gy1))
                    gx2 = max(0, min(top_w, gx2))
                    gy2 = max(0, min(top_h, gy2))
                    if gx2 <= gx1 or gy2 <= gy1:
                        continue

                    px_bbox = [gx1, gy1, gx2, gy2]

                    # 文本块：直接作为 text 元素
                    if blk_type in ["title", "text"]:
                        fig_mask.append(
                            {
                                "type": "text",
                                "bbox": px_bbox,
                                "text": text,
                                "text_level": 1 if blk_type == "title" else None,
                                "page_idx": 0,
                            }
                        )
                        text_count += 1
                    else:
                        # 非文本块：从顶层图再次裁剪成小图，作为 image 元素
                        try:
                            crop = top_img.crop((gx1, gy1, gx2, gy2))
                            icon_path = icons_raw_dir / f"blk_sub_{layout_idx}_{blk_idx}.png"
                            crop.save(icon_path)
                            fig_mask.append(
                                {
                                    "type": "image",
                                    "bbox": px_bbox,
                                    "img_path": str(icon_path),
                                    "page_idx": 0,
                                }
                            )
                            icon_count += 1
                        except Exception as e:
                            log.error(
                                f"[figure_mask] 子块裁剪失败 layout_idx={layout_idx}, blk_idx={blk_idx}, bbox={px_bbox}: {e}"
                            )
                            # 兜底：退化为文本元素，保持兼容
                            fig_mask.append(
                                {
                                    "type": "text",
                                    "bbox": px_bbox,
                                    "text": text,
                                    "text_level": None,
                                    "page_idx": 0,
                                }
                            )
                            text_count += 1
        else:
            # 正常路径：MinerU 输出多个元素，仍按原逻辑基于 fig_draft_path 裁剪
            for idx, it in enumerate(mineru_items):
                elem_type_raw = it.get("type") or ""
                elem_type = elem_type_raw.lower()
                bbox = it.get("bbox")
                text = (it.get("text") or it.get("content") or "").strip()

                if not bbox or len(bbox) != 4:
                    continue

                # 归一化 -> 像素坐标（基于原始 fig_ 图尺寸）
                x1n, y1n, x2n, y2n = bbox
                x1 = int(round(x1n * top_w))
                y1 = int(round(y1n * top_h))
                x2 = int(round(x2n * top_w))
                y2 = int(round(y2n * top_h))

                if x2 <= x1 or y2 <= y1:
                    continue

                px_bbox = [x1, y1, x2, y2]

                # 1) 只要有文字内容，一律作为文本元素
                if text:
                    fig_mask.append(
                        {
                            "type": "text",
                            "bbox": px_bbox,
                            "text": text,
                            "text_level": 1 if elem_type == "title" else None,
                            "page_idx": 0,
                        }
                    )
                    text_count += 1
                    continue

                # 2) 没有任何文字内容的块：一律裁图，当作 image，用于 icon / 元素图层
                try:
                    crop = top_img.crop((x1, y1, x2, y2))
                    icon_path = icons_raw_dir / f"blk_{idx}.png"
                    crop.save(icon_path)
                    icon_abs = str(icon_path)
                    fig_mask.append(
                        {
                            "type": "image",
                            "bbox": px_bbox,
                            "img_path": icon_abs,
                            "page_idx": 0,
                        }
                    )
                    icon_count += 1
                except Exception as e:
                    log.error(f"[figure_mask] 裁剪子图失败 idx={idx}, bbox={px_bbox}: {e}")
                    # 兜底：作为普通文本
                    fig_mask.append(
                        {
                            "type": "text",
                            "bbox": px_bbox,
                            "text": text,
                            "text_level": None,
                            "page_idx": 0,
                        }
                    )
                    text_count += 1

        type_counter = {}
        for e in fig_mask:
            t = e.get("type")
            type_counter[t] = type_counter.get(t, 0) + 1

        log.info(
            f"[figure_mask] fig_mask size = {len(fig_mask)}, "
            f"type distribution = {type_counter}, "
            f"title_text={text_count}, icons(raw)={icon_count}"
        )

        # 更新 state 的 fig_mask 信息
        state.fig_mask = fig_mask
        log.info(f"[figure_mask] 共解析出 {len(fig_mask)} 个元素 (via MinerU HTTP + SAM fallback, pixel bbox + raw icons)")

        return state
    
    async def figure_icon_bg_remover_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        把Mask里面的图标去除背景
        """
        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "icons"
        icons_dir.mkdir(parents=True, exist_ok=True)

        img_cnt = 0
        for item in state.fig_mask:
            if item.get('type') in ['image', 'table']:
                img_cnt += 1
                output_path = local_tool_for_bg_remove({
                    "image_path": item.get('img_path'),
                    "model_path": state.request.bg_rm_model,
                    "output_dir": str(icons_dir)
                })
                if output_path:
                    item['img_path'] = output_path
                    log.info(f"[figure_icon_bg_remover] background removed: {output_path}")
                else:
                    log.warning(f"[figure_icon_bg_remover] bg remove failed for {item.get('img_path')}")
        log.info(f"[figure_icon_bg_remover] processed image/table elements: {img_cnt}")

        # 抠图完成后，显式释放 RGB2.0 模型占用的显存
        try:
            free_bg_rm_model(model_path=state.request.bg_rm_model)
            log.info("[figure_icon_bg_remover] freed RMBG-2.0 model from GPU")
        except Exception as e:
            log.error(f"[figure_icon_bg_remover] free_bg_rm_model failed: {e}")

        return state

    async def figure_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成单页 PPT：
        - 第 1 页：原始组合页（layout EMF + MinerU 文本 + 图像）
        - 第 2 页：仅渲染所有 layout_items 的 EMF，用于检查 SAM 框架
        - 第 3 页：直接铺满整页的 full 内容 PNG（state.fig_draft_path）

        关键点：
        - layout_items 在 figure_layout_sam_node 中已经给出了像素坐标 bbox_px；
        - 这里完全沿用 add_text_element / add_image_element 使用的“像素 → 英寸 → Emu”规则，
          确保 EMF 背景框和内容层在同一几何坐标系下，避免因为单位不一致导致 EMF 肉眼不可见。
        - full 内容 PNG 页直接使用 fig_draft_path，按原图尺寸换算为英寸铺满整页。
        """
        try:
            # 从state获取输出目录（若未设置则自动初始化 outputs/paper2figure/<timestamp>）
            output_dir = Path(_ensure_result_path(state))
            output_dir.mkdir(parents=True, exist_ok=True)

            # 生成唯一文件名
            timestamp = int(time.time())
            ppt_filename = f"presentation_{timestamp}.pptx"
            ppt_path = output_dir / ppt_filename

            # 创建Presentation对象
            prs = Presentation()

            # 设置PPT尺寸，依据原始带内容图
            img = Image.open(state.fig_draft_path)
            width_px, height_px = img.size

            # --- 检查尺寸并计算缩放 (PPT限制56英寸) ---
            max_ppt_inches = 56.0
            dpi = 96.0
            max_pixels = int(max_ppt_inches * dpi)
            
            scale_ratio = 1.0
            if width_px > max_pixels or height_px > max_pixels:
                scale_ratio = max_pixels / max(width_px, height_px)
                # 留点余量，避免临界值误差
                scale_ratio *= 0.99 
                log.warning(f"[figure_ppt_generation] Image size ({width_px}x{height_px}) exceeds PPT limit. Scaling by {scale_ratio:.4f}")
                
                width_px = int(width_px * scale_ratio)
                height_px = int(height_px * scale_ratio)

            slide_width_px, slide_height_px = setup_presentation_size(prs, width_px, height_px)

            # 空白布局
            blank_slide_layout = prs.slide_layouts[6]

            def _add_layout_emf(slide, item) -> bool:
                """
                将 layout_item 中的 EMF 按像素 bbox 放到 slide 上，坐标逻辑与 add_image_element 保持一致。
                返回是否成功绘制。
                """
                emf_path = item.get("emf_path")
                if not emf_path or not os.path.exists(emf_path):
                    if emf_path:
                        log.warning(f"[figure_ppt_generation] emf_path 不存在: {emf_path}")
                    return False

                # 优先使用像素 bbox_px，其次退回原始 bbox（假定已是像素坐标）
                bbox = item.get("bbox_px") or item.get("bbox")
                if not bbox or len(bbox) != 4:
                    log.warning(f"[figure_ppt_generation] layout_item 缺少有效 bbox: {item}")
                    return False

                x1, y1, x2, y2 = bbox

                # 应用缩放
                if scale_ratio != 1.0:
                    x1 = int(x1 * scale_ratio)
                    y1 = int(y1 * scale_ratio)
                    x2 = int(x2 * scale_ratio)
                    y2 = int(y2 * scale_ratio)

                if x2 <= x1 or y2 <= y1:
                    log.warning(f"[figure_ppt_generation] 非法 bbox 像素坐标: {bbox} (scaled)")
                    return False

                # 像素 → 英寸，完全沿用 utils.pixels_to_inches 的规则
                left_in = pixels_to_inches(x1)
                top_in = pixels_to_inches(y1)
                width_in = pixels_to_inches(x2 - x1)
                height_in = pixels_to_inches(y2 - y1)

                log.info(f"[figure_ppt_generation] 添加 EMF：")
                log.info(f"  bbox 像素: [{x1}, {y1}, {x2}, {y2}]")
                log.info(f"  英寸坐标: left={left_in:.2f}, top={top_in:.2f}, width={width_in:.2f}, height={height_in:.2f}")
                log.info(f"  emf_path: {emf_path}")

                try:
                    slide.shapes.add_picture(
                        emf_path,
                        Inches(left_in),
                        Inches(top_in),
                        Inches(width_in),
                        Inches(height_in),
                    )
                    return True
                except Exception as e:
                    log.error(f"[figure_ppt_generation] add_picture EMF 失败: {emf_path}, {e}")
                    return False

            # =========================
            # 第 1 页：完整组合页
            # =========================
            slide_main = prs.slides.add_slide(blank_slide_layout)

            # 白色背景
            background = slide_main.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            # 1) 先渲染 layout_items (SAM + SVG + EMF 背景层)
            layout_drawn = 0
            for item in state.layout_items or []:
                if _add_layout_emf(slide_main, item):
                    layout_drawn += 1

            # 2) 再渲染 MinerU fig_mask（内容层）
            img_drawn = 0
            text_drawn = 0
            for element in state.fig_mask or []:
                # 应用缩放 (使用副本以免修改原数据)
                element_copy = element.copy()
                if scale_ratio != 1.0:
                    old_bbox = element_copy.get('bbox', [0,0,0,0])
                    if len(old_bbox) == 4:
                        element_copy['bbox'] = [
                            int(val * scale_ratio) for val in old_bbox
                        ]

                elem_type = element_copy.get('type', '')

                if elem_type == 'text':
                    add_text_element(slide_main, element_copy)
                    text_drawn += 1
                elif elem_type in ['image', 'table']:
                    add_image_element(slide_main, element_copy)
                    img_drawn += 1

            # =========================
            # 第 2 页：仅 EMF 调试页
            # =========================
            slide_emf = prs.slides.add_slide(blank_slide_layout)
            bg2 = slide_emf.background
            fill2 = bg2.fill
            fill2.solid()
            fill2.fore_color.rgb = RGBColor(255, 255, 255)

            layout_debug_drawn = 0
            for item in state.layout_items or []:
                if _add_layout_emf(slide_emf, item):
                    layout_debug_drawn += 1

            # =========================
            # 第 3 页：full 内容 PNG 原图页（铺满整页）
            # =========================
            slide_full = prs.slides.add_slide(blank_slide_layout)
            bg3 = slide_full.background
            fill3 = bg3.fill
            fill3.solid()
            fill3.fore_color.rgb = RGBColor(255, 255, 255)

            try:
                # 整幅图从 (0,0) 开始铺满整页，坐标/尺寸仍按像素->英寸转换
                left_in = pixels_to_inches(0)
                top_in = pixels_to_inches(0)
                width_in = pixels_to_inches(width_px)
                height_in = pixels_to_inches(height_px)

                slide_full.shapes.add_picture(
                    state.fig_draft_path,
                    Inches(left_in),
                    Inches(top_in),
                    Inches(width_in),
                    Inches(height_in),
                )
            except Exception as e:
                log.error(f"[figure_ppt_generation] add full PNG on page 3 failed: {e}")

            # 保存PPT
            prs.save(str(ppt_path))
            state.ppt_path = ppt_path
            print(f"PPT generated successfully: {ppt_path}")
            print(f"Slide size: {slide_width_px}x{slide_height_px} pixels")
            print(f"[MAIN] Total layout items: {len(state.layout_items)}, drawn: {layout_drawn}")
            print(f"[MAIN] Total content elements added: {len(state.fig_mask)}, text_drawn={text_drawn}, img_drawn={img_drawn}")
            print(f"[EMF_ONLY] layout items drawn: {layout_debug_drawn}")

        except Exception as e:
            print(f"Error generating PPT: {e}")

        return state

    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================
    def set_entry_node(state: Paper2FigureState) -> str:
        if(state.request.input_type == "PDF"):
            log.critical(f'进入PDF node ......')
            return "paper_idea_extractor"
        elif(state.request.input_type == "TEXT"):
            log.critical(f'进入TEXT node ......')
            return "figure_desc_generator"
        elif(state.request.input_type == "FIGURE"):
            log.critical(f'进入FIGURE node ......')
            return "figure_layout_sam"
        else:
            log.error(f"Invalid input type: {state.request.input_type}")
            return "_end_"

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        """
        _start_ 节点：确保本次 workflow 有一个统一的 result_path 根目录。
        - 若用户已在 state.result_path 传入自定义目录，则直接使用该目录；
        - 若未传入，则初始化为 get_project_root()/outputs/paper2figure/<timestamp>。
        """
        _ensure_result_path(state)
        return state

    nodes = {
        '_start_': _init_result_path,
        "paper_idea_extractor": paper_idea_extractor_node,
        "figure_desc_generator": figure_desc_generator_node,
        "figure_generator": figure_generator_node,
        "figure_layout_sam": figure_layout_sam_node,
        "figure_mask_generator": figure_mask_generator_node,
        "figure_icon_bg_remover": figure_icon_bg_remover_node,
        "figure_ppt_generator": figure_ppt_generation_node,
        '_end_': lambda state: state,  # 终止节点
    }

    # ------------------------------------------------------------------
    # EDGES  (从节点 A 指向节点 B)
    # ------------------------------------------------------------------
    edges = [
        ("paper_idea_extractor", "figure_desc_generator"),
        ("figure_desc_generator", "figure_generator"),
        ("figure_generator", "figure_layout_sam"),
        ("figure_layout_sam", "figure_mask_generator"),
        ("figure_mask_generator", "figure_icon_bg_remover"),
        ("figure_icon_bg_remover", "figure_ppt_generator"),
        ("figure_ppt_generator", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", set_entry_node)
    return builder
