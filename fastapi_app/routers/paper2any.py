from __future__ import annotations

import asyncio
from datetime import datetime
import uuid
from pathlib import Path
from typing import Optional
import httpx

from fastapi import APIRouter, File, Form, HTTPException, UploadFile, Request, Body
from fastapi.responses import FileResponse
from fastapi_app.routers.paper2video import paper2video_endpoint, FeaturePaper2VideoRequest, FeaturePaper2VideoResponse

from fastapi_app.schemas import Paper2FigureRequest, Paper2FigureResponse, VerifyLlmRequest, VerifyLlmResponse
from fastapi_app.workflow_adapters import run_paper2figure_wf_api
from dataflow_agent.utils import get_project_root
from fastapi_app.utils import _to_outputs_url, validate_invite_code  # noqa: F401
from dataflow_agent.logger import get_logger

log = get_logger(__name__)


# 全局信号量：控制重任务并发度（排队机制）
# 目前设为 1，即串行执行；如需并行可调大此值
task_semaphore = asyncio.Semaphore(1)

# 输出根目录：按任务类型 / 时间戳+UUID 组织
BASE_OUTPUT_DIR = Path("outputs")
PROJECT_ROOT = get_project_root()


router = APIRouter()


def create_run_dir(task_type: str) -> Path:
    """
    为一次请求创建独立目录：
        outputs/{task_type}/{timestamp}_{short_uuid}/
    并在其中创建 input/ 与 output/ 子目录。
    """
    ts = datetime.utcnow().strftime("%Y%m%dT%H%M%S")
    rid = uuid.uuid4().hex[:6]
    run_dir = BASE_OUTPUT_DIR / task_type / f"{ts}_{rid}"

    (run_dir / "input").mkdir(parents=True, exist_ok=True)
    (run_dir / "output").mkdir(parents=True, exist_ok=True)

    return run_dir


def create_dummy_pptx(output_path: Path, title: str, content: str) -> None:
    """
    生成一个非常简单的 PPTX 文件，作为占位 / Demo。

    若未来想使用真正的 workflow，可以在此处替换为实际调用逻辑。
    """
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        # 如果未安装 python-pptx，则写入一个占位的二进制文件，保证前端能正常下载
        output_path.write_bytes(b"Dummy PPTX content - please install python-pptx for real PPTX generation.")
        return

    prs = Presentation()
    # 使用标题 + 内容布局（通常为 1）
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

    prs.save(output_path)


@router.post("/verify-llm", response_model=VerifyLlmResponse)
async def verify_llm_connection(req: VerifyLlmRequest = Body(...)):
    """
    Verify LLM connection by sending a simple 'Hi' message from the backend.
    This avoids Mixed Content issues when the frontend is HTTPS and the LLM API is HTTP.
    """
    api_url = req.api_url.rstrip("/")
    if api_url.endswith("/chat/completions"):
        target_url = api_url
    else:
        target_url = f"{api_url}/chat/completions"

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {req.api_key}",
    }
    
    payload = {
        "model": req.model,
        "messages": [{"role": "user", "content": "Hi"}],
        "max_tokens": 1024
    }

    try:
        async with httpx.AsyncClient(timeout=10.0) as client:
            resp = await client.post(target_url, json=payload, headers=headers)
            
            if resp.status_code != 200:
                error_msg = f"API Error {resp.status_code}: {resp.text[:200]}"
                return VerifyLlmResponse(success=False, error=error_msg)
            
            return VerifyLlmResponse(success=True)
            
    except Exception as e:
        log.error(f"LLM Verification failed: {e}")
        return VerifyLlmResponse(success=False, error=str(e))


@router.get("/paper2figure/history_files")
async def list_paper2figure_history_files(
    request: Request,
    invite_code: str,
):
    """
    根据邀请码，列出该邀请码下面二级子目录中的所有历史输出文件（pptx/png/svg），
    即：outputs/{invite_code}/*/*.{pptx,png,svg} 以及 outputs/{invite_code}/*/*/{files} 这一层，
    不再往更深层递归。返回 URL 列表，前端可直接打开/下载。
    """
    # 邀请码校验
    # validate_invite_code(invite_code)

    project_root = get_project_root()
    base_dir = project_root / "outputs" / invite_code

    if not base_dir.exists():
        return {
            "success": True,
            "files": [],
        }

    file_urls: list[str] = []

    # 第一层：invite_code/level1
    for level1 in base_dir.iterdir():
        if not level1.is_dir():
            continue

        # 第二层：invite_code/level1/level2
        for level2 in level1.iterdir():
            if level2.is_file():
                # 若 level2 直接是文件，也纳入
                if level2.suffix.lower() in {".pptx", ".png", ".svg"}:
                    file_urls.append(_to_outputs_url(str(level2), request))
            elif level2.is_dir():
                # 只取该目录里的直接文件，不再往下递归
                for p in level2.iterdir():
                    if p.is_file() and p.suffix.lower() in {".pptx", ".png", ".svg"}:
                        file_urls.append(_to_outputs_url(str(p), request))

    # 排序：按路径字符串倒序，粗略实现“新文件在前”
    file_urls.sort(reverse=True)

    return {
        "success": True,
        "files": file_urls,
    }


@router.post("/paper2figure/generate")
async def generate_paper2figure(
    img_gen_model_name: str = Form(...),
    chat_api_url: str = Form(...),
    api_key: str = Form(...),
    input_type: str = Form(...),
    invite_code: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    file_kind: Optional[str] = Form(None),
    text: Optional[str] = Form(None),
    graph_type: str = Form("model_arch"),   # 'model_arch' | 'tech_route' | 'exp_data'
    language: str = Form("zh"),
    figure_complex: str = Form("easy"),
    style: str = Form("cartoon"),
):
    """
    Paper2Graph 接口（带邀请码校验 + workflow 调用）：

    - 前端通过 FormData 传入：
        - img_gen_model_name, chat_api_url, api_key, input_type, invite_code
        - file, file_kind, text
    - 路由层负责：
        - 校验邀请码
        - 保存文件到本地 input 目录
        - 将前端自定义 input_type/file_kind 映射为 Paper2FigureRequest 语义：
            - "file" + "pdf"   -> input_type="PDF",    input_content=PDF路径
            - "file" + "image" -> input_type="FIGURE", input_content=图片路径
            - "image"          -> input_type="FIGURE", input_content=图片路径
            - "text"           -> input_type="TEXT",   input_content=文本内容
        - 调用 run_paper2figure_wf_api
        - 返回生成的 PPTX 文件
    """
    # 0. 邀请码校验
    # validate_invite_code(invite_code)

    # 1. 基础参数校验
    if input_type in ("file", "image"):
        if file is None:
            raise HTTPException(
                status_code=400,
                detail="file is required when input_type is 'file' or 'image'",
            )
        if file_kind not in ("pdf", "image"):
            raise HTTPException(
                status_code=400,
                detail="file_kind must be 'pdf' or 'image'",
            )
    elif input_type == "text":
        if not text:
            raise HTTPException(
                status_code=400,
                detail="text is required when input_type is 'text'",
            )
    else:
        raise HTTPException(
            status_code=400,
            detail="invalid input_type, must be one of: file, text, image",
        )

    # 2. 创建本次请求的独立目录（按 graph_type 区分）并规范化难度
    if graph_type == "model_arch":
        task_type = "paper2fig"
        final_figure_complex = figure_complex or "easy"
    elif graph_type == "tech_route":
        task_type = "paper2tec"
        final_figure_complex = "easy"
    elif graph_type == "exp_data":
        task_type = "paper2exp"
        final_figure_complex = "easy"
    else:
        raise HTTPException(status_code=400, detail="invalid graph_type")

    # 语言：使用传入值或默认 zh（由 Form 默认保证）
    final_language = language

    run_dir = create_run_dir(task_type)
    input_dir = run_dir / "input"
    output_dir = run_dir / "output"

    # 3. 保存输入内容到 input/ 目录
    if input_type in ("file", "image"):
        original_name = file.filename or "uploaded"
        ext = Path(original_name).suffix or ""
        input_path = input_dir / f"input{ext}"
        content_bytes = await file.read()
        input_path.write_bytes(content_bytes)
    else:
        input_path = input_dir / "input.txt"
        input_path.write_text(text or "", encoding="utf-8")

    # 4. 将前端的 input_type/file_kind 映射为 Paper2FigureRequest 的语义
    if input_type in ("file", "image"):
        if file_kind == "pdf":
            real_input_type = "PDF"
            real_input_content = str(input_path)
        else:
            real_input_type = "FIGURE"
            real_input_content = str(input_path)
    elif input_type == "text":
        real_input_type = "TEXT"
        real_input_content = text or ""
    else:
        raise HTTPException(status_code=400, detail="unsupported input_type")

    # 5. 构造 Paper2FigureRequest
    p2f_req = Paper2FigureRequest(
        language=final_language,
        chat_api_url=chat_api_url,
        chat_api_key=api_key,
        api_key=api_key,
        # model 默认为 gpt-4o；如需前端控制可加字段
        model="gpt-4o",
        gen_fig_model=img_gen_model_name,
        input_type=real_input_type,        # "PDF" / "TEXT" / "FIGURE"
        input_content=real_input_content,  # 文件路径或文本
        aspect_ratio="16:9",
        graph_type=graph_type,
        style=style,
        figure_complex=final_figure_complex,
        invite_code=invite_code or "",
    )

    # 6. 重任务段：受信号量保护，调用真实 workflow
    async with task_semaphore:
        p2f_resp = await run_paper2figure_wf_api(p2f_req)

    # 7. 从 workflow 返回的路径读取 PPTX，并返回给前端
    raw_path = Path(p2f_resp.ppt_filename)

    # 若为相对路径，则以项目根目录为基准，避免工作目录变化导致找不到文件
    if not raw_path.is_absolute():
        ppt_path = PROJECT_ROOT / raw_path
    else:
        ppt_path = raw_path

    # 同时检查“存在 + 是文件”
    if not ppt_path.exists() or not ppt_path.is_file():
        raise HTTPException(
            status_code=500,
            detail=f"generated ppt file not found or not a file: {ppt_path}",
        )

    return FileResponse(
        path=str(ppt_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=ppt_path.name,
    )


@router.post("/paper2figure/generate_json", response_model=Paper2FigureResponse)
async def generate_paper2figure_json(
    request: Request,
    img_gen_model_name: str = Form(...),
    chat_api_url: str = Form(...),
    api_key: str = Form(...),
    input_type: str = Form(...),  # 'file' | 'text' | 'image'
    invite_code: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    file_kind: Optional[str] = Form(None),  # 'pdf' | 'image'
    text: Optional[str] = Form(None),
    graph_type: str = Form("model_arch"),   # 'model_arch' | 'tech_route' | 'exp_data'
    language: str = Form("zh"),
    style: str = Form("cartoon"),
):
    """
    Paper2Graph JSON 接口：

    - 与 /paper2figure/generate 使用同一套 FormData 参数
    - 但直接返回 Paper2FigureResponse(JSON)，包含:
        - ppt_filename
        - svg_filename
        - svg_image_filename
    """
    # 0. 邀请码校验
    # validate_invite_code(invite_code)

    # 1. 基础参数校验（与 generate_paper2figure 保持一致）
    if input_type in ("file", "image"):
        if file is None:
            raise HTTPException(
                status_code=400,
                detail="file is required when input_type is 'file' or 'image'",
            )
        if file_kind not in ("pdf", "image"):
            raise HTTPException(
                status_code=400,
                detail="file_kind must be 'pdf' or 'image'",
            )
    elif input_type == "text":
        if not text:
            raise HTTPException(
                status_code=400,
                detail="text is required when input_type is 'text'",
            )
    else:
        raise HTTPException(
            status_code=400,
            detail="invalid input_type, must be one of: file, text, image",
        )

    # 2. 创建本次请求的独立目录（按 graph_type 区分）
    if graph_type == "model_arch":
        task_type = "paper2fig"
    elif graph_type == "tech_route":
        task_type = "paper2tec"
    elif graph_type == "exp_data":
        task_type = "paper2exp"
    else:
        raise HTTPException(status_code=400, detail="invalid graph_type")

    run_dir = create_run_dir(task_type)
    input_dir = run_dir / "input"

    # 3. 保存输入内容到 input/ 目录
    if input_type in ("file", "image"):
        original_name = file.filename or "uploaded"
        ext = Path(original_name).suffix or ""
        input_path = input_dir / f"input{ext}"
        content_bytes = await file.read()
        input_path.write_bytes(content_bytes)
    else:
        input_path = input_dir / "input.txt"
        input_path.write_text(text or "", encoding="utf-8")

    # 4. 将前端的 input_type/file_kind 映射为 Paper2FigureRequest 的语义
    if input_type in ("file", "image"):
        if file_kind == "pdf":
            real_input_type = "PDF"
            real_input_content = str(input_path)
        else:
            real_input_type = "FIGURE"
            real_input_content = str(input_path)
    elif input_type == "text":
        real_input_type = "TEXT"
        real_input_content = text or ""
    else:
        raise HTTPException(status_code=400, detail="unsupported input_type")

    # 5. 构造 Paper2FigureRequest
    p2f_req = Paper2FigureRequest(
        language=language,
        chat_api_url=chat_api_url,
        chat_api_key=api_key,
        api_key=api_key,
        model="gpt-4o",
        gen_fig_model=img_gen_model_name,
        input_type=real_input_type,
        input_content=real_input_content,
        aspect_ratio="16:9",
        graph_type=graph_type,
        style=style,
        invite_code=invite_code or "",
    )

    # 6. 重任务段：受信号量保护，调用真实 workflow
    async with task_semaphore:
        p2f_resp = await run_paper2figure_wf_api(p2f_req)

    print(f"paper2figure response: {p2f_resp}")

    # 将绝对路径转换为前端可访问的完整 URL（包含协议、域名和端口）
    safe_ppt = _to_outputs_url(p2f_resp.ppt_filename, request)
    safe_svg = _to_outputs_url(p2f_resp.svg_filename, request) if p2f_resp.svg_filename else ""
    safe_png = _to_outputs_url(p2f_resp.svg_image_filename, request) if p2f_resp.svg_image_filename else ""

    # 新增：将本次任务输出目录下s所有相关文件路径转换为 URL
    safe_all_files: list[str] = []
    for abs_path in getattr(p2f_resp, "all_output_files", []) or []:
        if abs_path:
            safe_all_files.append(_to_outputs_url(abs_path, request))

    return Paper2FigureResponse(
        success=p2f_resp.success,
        ppt_filename=safe_ppt,
        svg_filename=safe_svg,
        svg_image_filename=safe_png,
        all_output_files=safe_all_files,
    )


@router.post("/paper2beamer/generate")
async def generate_paper2beamer(
    model_name: str = Form(...),
    chat_api_url: str = Form(...),
    api_key: str = Form(...),
    input_type: str = Form(...),  # 当前前端固定为 'file'
    invite_code: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    file_kind: Optional[str] = Form(None),  # 当前前端固定为 'pdf'
    language: str = Form(...),
):
    """
    paper2beamer 假接口（带邀请码校验）：

    - 需要前端在 FormData 中传入 invite_code，并在本地白名单文件中验证；
    - 接收前端上传的 PDF；
    - 为每次请求在 outputs/paper2beamer 下创建独立目录；
    - 使用全局信号量控制重任务串行执行；
    - 返回一个简单的 PPTX 文件，供前端下载测试。
    """
    # 0. 邀请码校验
    # validate_invite_code(invite_code)

    if input_type != "file":
        raise HTTPException(status_code=400, detail="paper2beamer currently only supports input_type='file'")

    if file is None:
        raise HTTPException(status_code=400, detail="file is required for paper2beamer")

    if file_kind not in ("pdf", None):
        # 允许 None（前端若未传），否则校验必须为 pdf
        raise HTTPException(status_code=400, detail="file_kind must be 'pdf' for paper2beamer")

    # 2. 创建本次请求的独立目录
    run_dir = create_run_dir("paper2beamer")
    input_dir = run_dir / "input"
    output_dir = run_dir / "output"

    # 3. 保存输入 PDF
    original_name = file.filename or "uploaded.pdf"
    ext = Path(original_name).suffix or ".pdf"
    input_path = input_dir / f"input{ext}"
    content_bytes = await file.read()
    input_path.write_bytes(content_bytes)
    abs_input_path = input_path.resolve()
    # saved_input_name = input_path.name

    # 4. 重任务段：受信号量保护，确保排队执行
    async with task_semaphore:
        # output_pptx = output_dir / "paper2beamer.pdf"
        # demo_title = "paper2beamer Demo"
        # content = (
        #     f"model_name: {model_name}\n"
        #     f"chat_api_url: {chat_api_url}\n"
        #     f"input_type: {input_type}\n"
        #     f"file_kind: {file_kind or 'pdf'}\n"
        #     f"saved_input: {saved_input_name}\n"
        # )
        # create_dummy_pptx(output_pptx, demo_title, demo_content)
        # create_pdf(output_pptx, demo_title, content)
        req = FeaturePaper2VideoRequest(
            model=model_name,
            chat_api_url=chat_api_url,
            api_key=api_key,
            pdf_path=str(abs_input_path),
            img_path="",
            language=language,
        )
        resp: FeaturePaper2VideoResponse = await paper2video_endpoint(req)
        if not resp.success:
            raise HTTPException(status_code=500, detail="Paper to PPT generation failed.")
        output_path = resp.ppt_path
        output_path = Path(output_path)

    # 5. 返回 PPTX 文件
    # return FileResponse(
    #     path=output_pptx,
    #     media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    #     filename="paper2beamer.pdf",
    # )
    return FileResponse(
    path=output_path,
    media_type="application/pdf",
    filename="paper2beamer.pdf",
)
